#!/usr/bin/env python3
"""
Fast Text Extraction Pipeline
=============================

Ultra-fast text extraction for text-heavy documents.
Target: 100+ pages/second for PDFs, 300+ files/minute for other formats.

This module bypasses heavy ML models for simple text extraction,
generating immediate results with placeholders for visual elements.
"""

import time
import tempfile
import subprocess
from pathlib import Path
from typing import Dict, List, Optional, Union, Any
from dataclasses import dataclass
from enum import Enum
import json
import hashlib

# PDF libraries (multiple options for maximum compatibility)
try:
    import pymupdf as fitz
    PYMUPDF_AVAILABLE = True
except ImportError:
    try:
        import fitz
        PYMUPDF_AVAILABLE = True
    except ImportError:
        PYMUPDF_AVAILABLE = False

try:
    from pypdfium2 import PdfDocument
    PYPDFIUM2_AVAILABLE = True
except ImportError:
    PYPDFIUM2_AVAILABLE = False

# Office document libraries
try:
    from docx import Document as DocxDocument
    PYTHON_DOCX_AVAILABLE = True
except ImportError:
    PYTHON_DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    PYTHON_PPTX_AVAILABLE = False

try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

# Web and markup
try:
    from bs4 import BeautifulSoup
    BS4_AVAILABLE = True
except ImportError:
    BS4_AVAILABLE = False

import csv
import pandas as pd


class ExtractionMethod(Enum):
    """Text extraction methods by speed preference."""
    PYMUPDF = "pymupdf"           # Fastest for PDFs
    PYPDFIUM2 = "pypdfium2"       # Fallback for PDFs
    PYTHON_DOCX = "python_docx"   # Office docs
    PYTHON_PPTX = "python_pptx"   # PowerPoint
    OPENPYXL = "openpyxl"         # Excel
    BEAUTIFULSOUP = "bs4"         # HTML
    BUILTIN = "builtin"           # Text files, CSV, etc.
    SUBPROCESS = "subprocess"     # External tools fallback


@dataclass
class ExtractionResult:
    """Result of fast text extraction."""
    success: bool
    text_content: str
    metadata: Dict[str, Any]
    visual_placeholders: List[Dict[str, Any]]
    extraction_method: ExtractionMethod
    extraction_time: float
    page_count: int
    word_count: int
    char_count: int
    error_message: Optional[str] = None


@dataclass
class VisualPlaceholder:
    """Placeholder for visual elements that need enhanced processing."""
    placeholder_id: str
    element_type: str              # "image", "table", "formula", "chart"
    page_number: int
    bbox: Optional[tuple] = None   # (x0, y0, x1, y1) if available
    description: str = "Visual element processing queued"
    priority: int = 3              # 1-5, higher = more important


class FastTextExtractor:
    """Ultra-fast text extraction for immediate results."""
    
    def __init__(self, log_to_file=False, text_only_mode=False):
        self.placeholder_counter = 0
        self.log_to_file = log_to_file
        self.log_file = None
        self.text_only_mode = text_only_mode  # Skip all image extraction when True
        
        if self.log_to_file:
            # Create output directory and log file
            from datetime import datetime
            output_dir = Path('/home/corey/projects/docling/cli/output/latest')
            output_dir.mkdir(parents=True, exist_ok=True)
            
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            self.log_file = output_dir / f"extraction_log_{timestamp}.txt"
            
            # Initialize log file
            with open(self.log_file, 'w') as f:
                f.write(f"=== FAST TEXT EXTRACTION LOG - {datetime.now()} ===\n\n")
    
    def _log(self, message: str):
        """Log message to both console and file if enabled."""
        print(message)
        if self.log_to_file and self.log_file:
            with open(self.log_file, 'a') as f:
                f.write(message + '\n')
        
    def extract(self, file_path: Path) -> ExtractionResult:
        """
        Extract text as quickly as possible from any supported document.
        
        Target performance:
        - PDFs: 100+ pages/second
        - Office docs: 300+ files/minute  
        - Text formats: 1000+ files/minute
        """
        start_time = time.time()
        
        if not file_path.exists():
            return ExtractionResult(
                success=False,
                text_content="",
                metadata={},
                visual_placeholders=[],
                extraction_method=ExtractionMethod.BUILTIN,
                extraction_time=0.0,
                page_count=0,
                word_count=0,
                char_count=0,
                error_message=f"File not found: {file_path}"
            )
        
        # Route by file extension
        suffix = file_path.suffix.lower()
        
        try:
            if suffix == '.pdf':
                result = self._extract_pdf(file_path)
            elif suffix in ['.docx', '.doc']:
                result = self._extract_docx(file_path)
            elif suffix in ['.pptx', '.ppt']:
                result = self._extract_pptx(file_path)
            elif suffix in ['.xlsx', '.xls']:
                result = self._extract_xlsx(file_path)
            elif suffix in ['.html', '.htm']:
                result = self._extract_html(file_path)
            elif suffix in ['.md', '.markdown']:
                result = self._extract_markdown(file_path)
            elif suffix == '.csv':
                result = self._extract_csv(file_path)
            elif suffix in ['.txt', '.text']:
                result = self._extract_text(file_path)
            elif suffix == '.json':
                result = self._extract_json(file_path)
            elif suffix in ['.xml']:
                result = self._extract_xml(file_path)
            else:
                # Fallback: try as text
                result = self._extract_text(file_path)
                
        except Exception as e:
            result = ExtractionResult(
                success=False,
                text_content="",
                metadata={},
                visual_placeholders=[],
                extraction_method=ExtractionMethod.BUILTIN,
                extraction_time=0.0,
                page_count=0,
                word_count=0,
                char_count=0,
                error_message=str(e)
            )
        
        # Finalize timing and stats
        result.extraction_time = time.time() - start_time
        if result.success and result.text_content:
            result.word_count = len(result.text_content.split())
            result.char_count = len(result.text_content)
        
        return result
    
    def _extract_pdf(self, file_path: Path) -> ExtractionResult:
        """Ultra-fast PDF text extraction."""
        
        if PYMUPDF_AVAILABLE:
            return self._extract_pdf_pymupdf(file_path)
        elif PYPDFIUM2_AVAILABLE:
            return self._extract_pdf_pypdfium2(file_path)
        else:
            # Fallback to subprocess (pdftotext if available)
            return self._extract_pdf_subprocess(file_path)
    
    def _extract_pdf_pymupdf(self, file_path: Path) -> ExtractionResult:
        """PyMuPDF extraction with pre-detection of mathematical regions."""
        doc = fitz.open(str(file_path))
        
        text_content = []
        visual_placeholders = []
        page_count = len(doc)
        
        for page_num in range(page_count):
            page = doc[page_num]
            
            if self.text_only_mode:
                # FAST PATH: Skip all detection in text-only mode
                page_text = page.get_text("text", sort=True)
                text_content.append(page_text)
            else:
                # FULL PATH: Detect and preserve special regions
                # First, detect mathematical and table regions BEFORE text extraction
                math_regions, table_regions, figure_regions = self._detect_special_regions(page, page_num + 1)
                
                # Extract text blocks with preserved regions
                page_text = self._extract_text_with_preservation(page, math_regions, table_regions, figure_regions, page_num + 1)
                
                # Create placeholders for all detected regions
                page_placeholders = math_regions + table_regions + figure_regions
                
                text_content.append(page_text)
                visual_placeholders.extend(page_placeholders)
        
        doc.close()
        
        full_text = "\n\n".join(text_content)
        
        return ExtractionResult(
            success=True,
            text_content=full_text,
            metadata={
                "file_path": str(file_path),
                "file_size": file_path.stat().st_size,
                "pdf_version": "unknown"
            },
            visual_placeholders=visual_placeholders,
            extraction_method=ExtractionMethod.PYMUPDF,
            extraction_time=0.0,  # Will be set by caller
            page_count=page_count,
            word_count=0,  # Will be calculated
            char_count=0   # Will be calculated
        )
    
    def _extract_pdf_pypdfium2(self, file_path: Path) -> ExtractionResult:
        """Pypdfium2 extraction with visual element detection."""
        doc = PdfDocument(str(file_path))
        page_count = len(doc)
        
        text_content = []
        visual_placeholders = []
        
        for i in range(page_count):
            page = doc[i]
            textpage = page.get_textpage()
            page_text = textpage.get_text_range()
            text_content.append(page_text)
            
            # Skip all detection in text-only mode
            if self.text_only_mode:
                continue
            
            # Detect figure references in text (charts, maps, diagrams)
            import re
            figure_matches = re.findall(r'Figure\s+(\d+):\s*([^.]+)', page_text, re.IGNORECASE)
            for fig_num, fig_desc in figure_matches:
                # Determine element type based on description
                element_type = "image"  # default
                desc_lower = fig_desc.lower()
                if any(word in desc_lower for word in ["chart", "graph", "bar", "plot"]):
                    element_type = "chart"
                elif any(word in desc_lower for word in ["map", "location", "geographic"]):
                    element_type = "map"
                elif any(word in desc_lower for word in ["diagram", "schematic"]):
                    element_type = "diagram"
                
                placeholder = self._create_placeholder(
                    element_type,
                    i + 1,
                    f"Figure {fig_num}: {fig_desc.strip()}"
                )
                visual_placeholders.append(placeholder)
            
            # Quick table detection (heuristic)
            if "|" in page_text or "┌" in page_text or "├" in page_text:
                placeholder = self._create_placeholder(
                    "table",
                    i + 1, 
                    f"Table detected on page {i + 1}"
                )
                visual_placeholders.append(placeholder)
        
        full_text = "\n\n".join(text_content)
        
        return ExtractionResult(
            success=True,
            text_content=full_text,
            metadata={"file_path": str(file_path)},
            visual_placeholders=visual_placeholders,
            extraction_method=ExtractionMethod.PYPDFIUM2,
            extraction_time=0.0,
            page_count=page_count,
            word_count=0,
            char_count=0
        )
    
    def _extract_pdf_subprocess(self, file_path: Path) -> ExtractionResult:
        """Subprocess fallback (pdftotext, etc.)."""
        try:
            # Try pdftotext if available
            result = subprocess.run(
                ['pdftotext', str(file_path), '-'],
                capture_output=True,
                text=True,
                timeout=30
            )
            
            if result.returncode == 0:
                return ExtractionResult(
                    success=True,
                    text_content=result.stdout,
                    metadata={"file_path": str(file_path)},
                    visual_placeholders=[],
                    extraction_method=ExtractionMethod.SUBPROCESS,
                    extraction_time=0.0,
                    page_count=1,  # Can't easily determine
                    word_count=0,
                    char_count=0
                )
            else:
                raise subprocess.CalledProcessError(result.returncode, 'pdftotext')
                
        except (subprocess.CalledProcessError, FileNotFoundError):
            # pdftotext not available or failed
            return ExtractionResult(
                success=False,
                text_content="",
                metadata={},
                visual_placeholders=[],
                extraction_method=ExtractionMethod.SUBPROCESS,
                extraction_time=0.0,
                page_count=0,
                word_count=0,
                char_count=0,
                error_message="PDF extraction failed: no suitable library available"
            )
    
    def _extract_docx(self, file_path: Path) -> ExtractionResult:
        """Fast DOCX text extraction."""
        if not PYTHON_DOCX_AVAILABLE:
            return self._extraction_failed("python-docx not available")
        
        doc = DocxDocument(str(file_path))
        
        # Extract text from paragraphs
        text_content = []
        visual_placeholders = []
        
        for para in doc.paragraphs:
            text_content.append(para.text)
        
        # Extract text from tables (basic)
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = [cell.text for cell in row.cells]
                table_text.append(" | ".join(row_text))
            
            # Add as placeholder for enhanced processing
            placeholder = self._create_placeholder(
                "table",
                1,  # DOCX doesn't have page numbers in same way
                f"Table with {len(table.rows)} rows"
            )
            visual_placeholders.append(placeholder)
            
            text_content.append("\n".join(table_text))
        
        # Check for images
        from docx.document import Document as DocxDoc
        relationships = doc.part.rels
        image_count = sum(1 for r in relationships.values() 
                         if "image" in r.target_ref)
        
        if image_count > 0:
            placeholder = self._create_placeholder(
                "image",
                1,
                f"{image_count} images in document"
            )
            visual_placeholders.append(placeholder)
        
        full_text = "\n\n".join(text_content)
        
        return ExtractionResult(
            success=True,
            text_content=full_text,
            metadata={
                "file_path": str(file_path),
                "paragraph_count": len(doc.paragraphs),
                "table_count": len(doc.tables),
                "image_count": image_count
            },
            visual_placeholders=visual_placeholders,
            extraction_method=ExtractionMethod.PYTHON_DOCX,
            extraction_time=0.0,
            page_count=1,  # DOCX is continuous
            word_count=0,
            char_count=0
        )
    
    def _extract_pptx(self, file_path: Path) -> ExtractionResult:
        """Fast PPTX text extraction."""
        if not PYTHON_PPTX_AVAILABLE:
            return self._extraction_failed("python-pptx not available")
        
        prs = Presentation(str(file_path))
        
        text_content = []
        visual_placeholders = []
        slide_count = len(prs.slides)
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
                elif shape.shape_type == 13:  # Picture
                    placeholder = self._create_placeholder(
                        "image",
                        slide_num,
                        f"Image on slide {slide_num}"
                    )
                    visual_placeholders.append(placeholder)
            
            if slide_text:
                text_content.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_text))
        
        full_text = "\n\n".join(text_content)
        
        return ExtractionResult(
            success=True,
            text_content=full_text,
            metadata={
                "file_path": str(file_path),
                "slide_count": slide_count
            },
            visual_placeholders=visual_placeholders,
            extraction_method=ExtractionMethod.PYTHON_PPTX,
            extraction_time=0.0,
            page_count=slide_count,
            word_count=0,
            char_count=0
        )
    
    def _extract_xlsx(self, file_path: Path) -> ExtractionResult:
        """Fast Excel text extraction."""
        if not OPENPYXL_AVAILABLE:
            return self._extraction_failed("openpyxl not available")
        
        wb = openpyxl.load_workbook(str(file_path), data_only=True)
        
        text_content = []
        visual_placeholders = []
        
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            
            sheet_text = [f"--- Sheet: {sheet_name} ---"]
            
            # Extract cell values efficiently
            for row in ws.iter_rows(values_only=True):
                row_text = [str(cell) if cell is not None else "" for cell in row]
                if any(row_text):  # Skip empty rows
                    sheet_text.append(" | ".join(row_text))
            
            text_content.append("\n".join(sheet_text))
            
            # Add sheet as table placeholder
            placeholder = self._create_placeholder(
                "table",
                1,
                f"Spreadsheet: {sheet_name}"
            )
            visual_placeholders.append(placeholder)
        
        full_text = "\n\n".join(text_content)
        
        return ExtractionResult(
            success=True,
            text_content=full_text,
            metadata={
                "file_path": str(file_path),
                "sheet_count": len(wb.sheetnames),
                "sheet_names": wb.sheetnames
            },
            visual_placeholders=visual_placeholders,
            extraction_method=ExtractionMethod.OPENPYXL,
            extraction_time=0.0,
            page_count=len(wb.sheetnames),
            word_count=0,
            char_count=0
        )
    
    def _extract_html(self, file_path: Path) -> ExtractionResult:
        """Fast HTML text extraction."""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            html_content = f.read()
        
        if BS4_AVAILABLE:
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            text_content = soup.get_text()
            
            # Clean up whitespace
            lines = (line.strip() for line in text_content.splitlines())
            text_content = '\n'.join(line for line in lines if line)
            
            # Detect visual elements
            visual_placeholders = []
            images = soup.find_all('img')
            if images:
                placeholder = self._create_placeholder(
                    "image",
                    1,
                    f"{len(images)} images in HTML"
                )
                visual_placeholders.append(placeholder)
            
            tables = soup.find_all('table')
            if tables:
                placeholder = self._create_placeholder(
                    "table", 
                    1,
                    f"{len(tables)} tables in HTML"
                )
                visual_placeholders.append(placeholder)
                
        else:
            # Fallback: crude HTML stripping
            import re
            text_content = re.sub('<[^<]+?>', '', html_content)
            visual_placeholders = []
        
        return ExtractionResult(
            success=True,
            text_content=text_content,
            metadata={"file_path": str(file_path)},
            visual_placeholders=visual_placeholders,
            extraction_method=ExtractionMethod.BEAUTIFULSOUP if BS4_AVAILABLE else ExtractionMethod.BUILTIN,
            extraction_time=0.0,
            page_count=1,
            word_count=0,
            char_count=0
        )
    
    def _extract_markdown(self, file_path: Path) -> ExtractionResult:
        """Fast Markdown text extraction."""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text_content = f.read()
        
        # Markdown is already mostly text - minimal processing needed
        # Could add placeholder detection for ![image] syntax
        
        visual_placeholders = []
        import re
        
        # Find image references
        image_matches = re.findall(r'!\[([^\]]*)\]\([^)]+\)', text_content)
        if image_matches:
            placeholder = self._create_placeholder(
                "image",
                1,
                f"{len(image_matches)} image references"
            )
            visual_placeholders.append(placeholder)
        
        # Find table indicators (basic)
        if '|' in text_content and '---' in text_content:
            placeholder = self._create_placeholder(
                "table",
                1,
                "Markdown tables detected"
            )
            visual_placeholders.append(placeholder)
        
        return ExtractionResult(
            success=True,
            text_content=text_content,
            metadata={"file_path": str(file_path)},
            visual_placeholders=visual_placeholders,
            extraction_method=ExtractionMethod.BUILTIN,
            extraction_time=0.0,
            page_count=1,
            word_count=0,
            char_count=0
        )
    
    def integrate_visual_results(self, text_content: str, visual_jobs: List) -> str:
        """
        Integrate visual processing results back into the text content.
        
        Args:
            text_content: Original text with visual placeholders
            visual_jobs: List of completed visual jobs from queue manager
            
        Returns:
            Enhanced text with placeholders replaced by actual content
        """
        import re
        enhanced_text = text_content
        
        print(f"🔄 Integrating {len(visual_jobs)} visual results...")
        
        for job in visual_jobs:
            if hasattr(job, 'enhanced_content') and job.enhanced_content:
                # Use the stored placeholder_id if available, otherwise fall back to job_id logic
                if hasattr(job, 'placeholder_id') and job.placeholder_id:
                    placeholder_id = job.placeholder_id
                    print(f"   🏷️  Using placeholder_id: {placeholder_id}")
                elif hasattr(job, 'job_id'):
                    # Extract numeric ID from job_id if available (legacy fallback)
                    numeric_id = re.search(r'(\d+)', job.job_id)
                    if numeric_id:
                        placeholder_id = f"visual_{numeric_id.group(1).zfill(4)}"
                    else:
                        placeholder_id = f"visual_{job.page_number:04d}"
                else:
                    placeholder_id = f"visual_{job.page_number:04d}"
                
                print(f"   Processing {placeholder_id} ({job.element_type.value})")
                
                # Extract relevant content based on element type
                if job.element_type.value == 'table':
                    # Extract and clean markdown tables from enhanced content
                    table_content = self._extract_clean_table(job.enhanced_content)
                    
                    if table_content:
                        # Replace placeholder with actual table
                        placeholder_pattern = rf'`\[{placeholder_id}\]`'
                        if placeholder_pattern in enhanced_text:
                            enhanced_text = re.sub(re.escape(placeholder_pattern), table_content, enhanced_text)
                            print(f"   ✅ Replaced {placeholder_id} with table")
                        else:
                            print(f"   ⚠️  Placeholder {placeholder_pattern} not found in text")
                        continue

                elif job.element_type.value == 'formula':
                    # Extract and convert mathematical content to LaTeX/markdown
                    formula_content = self._extract_clean_formula(job.enhanced_content)
                    
                    if formula_content:
                        placeholder_pattern = rf'`\[{placeholder_id}\]`'
                        if placeholder_pattern in enhanced_text:
                            enhanced_text = re.sub(re.escape(placeholder_pattern), formula_content, enhanced_text)
                            print(f"   ✅ Replaced {placeholder_id} with formula")
                        else:
                            print(f"   ⚠️  Placeholder {placeholder_pattern} not found in text")
                        continue
                
                elif job.element_type.value == 'image':
                    # For images, extract the page-specific content from VLM output
                    print(f"   🔍 Processing {placeholder_id} ({job.element_type.value}) on page {job.page_number}")
                    print(f"   📄 VLM Content Length: {len(job.enhanced_content)} chars")
                    
                    # Check what type of content we have
                    if '|' in job.enhanced_content and job.enhanced_content.count('|') > 5:
                        print(f"   📊 MARKDOWN TABLE detected in content")
                    if 'chart' in job.enhanced_content.lower() or 'graph' in job.enhanced_content.lower():
                        print(f"   📈 CHART/GRAPH content detected")
                    
                    description = None
                    
                    # Try to extract page-specific content first
                    page_specific_content = self._extract_page_specific_content(job.enhanced_content, job.page_number)
                    if page_specific_content:
                        print(f"   📄 Found page-specific content: {repr(page_specific_content[:200])}")
                        
                        # Save page-specific content for verification
                        verification_file = f"/home/corey/projects/docling/cli/temp/page_{job.page_number}_content.md"
                        clean_full_content = self._remove_binary_data(job.enhanced_content)
                        with open(verification_file, 'w') as f:
                            f.write(f"# Page {job.page_number} Extracted Content (Binary Data Removed)\n\n")
                            f.write(f"**Original VLM Content Length**: {len(job.enhanced_content)} chars\n\n")
                            f.write(f"**Clean VLM Content Length**: {len(clean_full_content)} chars\n\n")
                            f.write(f"**Page-Specific Content Length**: {len(page_specific_content)} chars\n\n")
                            f.write("## Page-Specific Content:\n\n")
                            f.write(page_specific_content)
                            f.write(f"\n\n## Clean VLM Content (first 1000 chars):\n\n")
                            f.write(clean_full_content[:1000])
                        print(f"   💾 Saved clean page content to: {verification_file}")
                        
                        job.enhanced_content = page_specific_content  # Use page-specific content for pattern matching
                    
                    # First, try specific description patterns - prioritize tables and charts
                    desc_patterns = [
                        # Chart/table patterns (higher priority)
                        r'((?:\|[^\n]+\|[^\n]*\n\|[-\s:|]+\|[^\n]*\n(?:\|[^\n]+\|[^\n]*\n)*)+)',  # Full markdown tables
                        r'chart.*?shows?\s+([^.]+\.)',
                        r'graph.*?shows?\s+([^.]+\.)', 
                        r'map.*?shows?\s+([^.]+\.)',
                        r'proportions?\s+([^.]+\.)',
                        r'location\s+map\s+([^.]+\.)',
                        # Generic patterns (lower priority)
                        r'description:\s*([^\n]+)',
                        r'shows?\s+([^.]+\.)',
                        r'depicts?\s+([^.]+\.)',
                        r'figure\s+\d+[:.]\s*([^.]+\.)',
                        r'image\s+shows?\s+([^.]+\.)',
                        r'this\s+(?:image|figure|picture)\s+(?:shows?|depicts?|contains?)\s+([^.]+\.)'
                    ]
                    
                    for i, pattern in enumerate(desc_patterns):
                        print(f"   🔍 Trying pattern {i+1}: {pattern}")
                        desc_match = re.search(pattern, job.enhanced_content, re.IGNORECASE)
                        if desc_match:
                            description = desc_match.group(1).strip()
                            print(f"   ✅ Pattern {i+1} matched: {repr(description[:100])}")
                            break
                        else:
                            print(f"   ❌ Pattern {i+1} no match")
                    
                    # If no pattern matches, use the full enhanced content (cleaned up)
                    if not description and job.enhanced_content:
                        print(f"   🔄 No patterns matched, using full content as fallback")
                        # Remove binary data first, then clean up formatting
                        cleaned_content = self._remove_binary_data(job.enhanced_content)
                        cleaned_content = re.sub(r'[#*`\[\]]+', '', cleaned_content)
                        cleaned_content = ' '.join(cleaned_content.split())  # Normalize whitespace
                        print(f"   📝 Cleaned content: {repr(cleaned_content[:200])}")
                        if len(cleaned_content) > 20:  # Only use if substantial content
                            description = cleaned_content
                            print(f"   ✅ Using full content as description ({len(description)} chars)")
                    
                    if description:
                        placeholder_pattern = rf'`\[{placeholder_id}\]`'
                        print(f"   🔍 Looking for pattern: {placeholder_pattern}")
                        if placeholder_pattern in enhanced_text:
                            enhanced_text = re.sub(re.escape(placeholder_pattern), f"*{description}*", enhanced_text)
                            print(f"   ✅ Replaced {placeholder_id} with description: {description[:50]}...")
                        else:
                            print(f"   ❌ Pattern {placeholder_pattern} not found in text")
                            print(f"   📄 Text contains: {enhanced_text.count('visual_')} visual references")
                        continue
                
                # Fallback: replace placeholder with a generic message
                placeholder_pattern = rf'`\[{placeholder_id}\]`'
                if placeholder_pattern in enhanced_text:
                    fallback_text = f"*Enhanced {job.element_type.value} content processed*"
                    enhanced_text = re.sub(re.escape(placeholder_pattern), fallback_text, enhanced_text)
                    print(f"   ✅ Replaced {placeholder_id} with fallback")
                else:
                    print(f"   ⚠️  No placeholder found for {placeholder_id}")
            else:
                print(f"   ❌ Job has no enhanced content: {job}")
        
        # Only remove the "Visual Elements" section if all placeholders were successfully replaced
        if self._all_placeholders_replaced(enhanced_text):
            enhanced_text = self._remove_visual_elements_section(enhanced_text)
            print("   🧹 Removed Visual Elements section - all content integrated")
        else:
            print("   📋 Keeping Visual Elements section - some placeholders remain")
        
        return enhanced_text
    
    def _extract_clean_table(self, content: str) -> str:
        """Extract and clean table content from VLM output."""
        import re
        
        # Remove binary data first
        content = self._remove_binary_data(content)
        
        # Try different table patterns
        table_patterns = [
            r'(\|[^\n]+\|[^\n]*\n\|[-\s:|]+\|[^\n]*\n(?:\|[^\n]+\|[^\n]*\n)*)',  # Standard markdown tables
            r'(\|.*\|.*\n\|[-\s\|:]*\|.*\n(?:\|.*\|.*\n)*)',  # Alternative pattern
            r'((?:[^\n\|]*\|[^\n\|]*\|[^\n]*\n){3,})',  # Relaxed table pattern
        ]
        
        for pattern in table_patterns:
            tables = re.findall(pattern, content, re.MULTILINE)
            if tables:
                table = tables[0].strip()
                # Clean up the table
                lines = table.split('\n')
                cleaned_lines = []
                for line in lines:
                    # Ensure proper table formatting
                    if '|' in line:
                        line = line.strip()
                        if not line.startswith('|'):
                            line = '| ' + line
                        if not line.endswith('|'):
                            line = line + ' |'
                        cleaned_lines.append(line)
                
                if len(cleaned_lines) >= 3:  # Header + separator + at least one row
                    return '\n'.join(cleaned_lines)
        
        return None
    
    def _extract_clean_formula(self, content: str) -> str:
        """Extract and clean mathematical formula from VLM output."""
        import re
        
        # Remove binary data first  
        content = self._remove_binary_data(content)
        
        # First, check if we have preserved original mathematical text
        # Look for mathematical content that was preserved during extraction
        if hasattr(content, 'original_text'):
            original = content.original_text
            # Convert preserved mathematical symbols to LaTeX
            return self._convert_math_to_latex(original)
        
        # Look for LaTeX-style math expressions in VLM output
        latex_patterns = [
            r'\$\$([^$]+)\$\$',  # Display math
            r'\$([^$\n]+)\$',    # Inline math
            r'\\begin\{equation\}(.*?)\\end\{equation\}',  # Equation blocks
            r'\\begin\{align\}(.*?)\\end\{align\}',        # Align blocks
            r'\\[(.*?)\\]',      # LaTeX display math
        ]
        
        for pattern in latex_patterns:
            matches = re.findall(pattern, content, re.DOTALL)
            if matches:
                formula = matches[0].strip()
                # Clean up common artifacts
                formula = re.sub(r'\\text\{([^}]+)\}', r'\1', formula)  # Remove text wrappers
                formula = re.sub(r'\\mathrm\{([^}]+)\}', r'\1', formula)  # Remove mathrm wrappers
                return f"$${formula}$$"  # Return as display math
        
        # If no LaTeX found, convert mathematical symbols and expressions
        return self._convert_math_to_latex(content)
    
    def _convert_math_to_latex(self, text: str) -> str:
        """Convert mathematical text with Unicode symbols to LaTeX."""
        import re
        
        # Symbol mapping for common mathematical Unicode to LaTeX
        symbol_map = {
            '∑': r'\sum',
            '∫': r'\int', 
            '∏': r'\prod',
            '√': r'\sqrt',
            '±': r'\pm',
            '∞': r'\infty',
            '≈': r'\approx',
            '≤': r'\leq',
            '≥': r'\geq', 
            '≠': r'\neq',
            '∂': r'\partial',
            '∇': r'\nabla',
            '∀': r'\forall',
            '∃': r'\exists',
            '∈': r'\in',
            '∉': r'\notin',
            '⊂': r'\subset',
            '⊃': r'\supset',
            '∪': r'\cup',
            '∩': r'\cap',
            '∧': r'\land',
            '∨': r'\lor',
            '¬': r'\neg',
            '→': r'\to',
            '↔': r'\leftrightarrow',
            '∴': r'\therefore',
            '∵': r'\because',
            # Greek letters
            'α': r'\alpha', 'β': r'\beta', 'γ': r'\gamma', 'δ': r'\delta',
            'ε': r'\varepsilon', 'ζ': r'\zeta', 'η': r'\eta', 'θ': r'\theta',
            'λ': r'\lambda', 'μ': r'\mu', 'π': r'\pi', 'σ': r'\sigma',
            'φ': r'\phi', 'ψ': r'\psi', 'ω': r'\omega',
            'Θ': r'\Theta', 'Δ': r'\Delta', 'Φ': r'\Phi', 'Ψ': r'\Psi', 'Ω': r'\Omega'
        }
        
        # Replace Unicode symbols with LaTeX
        latex_text = text
        for unicode_sym, latex_sym in symbol_map.items():
            latex_text = latex_text.replace(unicode_sym, latex_sym)
        
        # Fix subscripts and superscripts
        latex_text = re.sub(r'([a-zA-Z])_([0-9]+)', r'\1_{\2}', latex_text)
        latex_text = re.sub(r'([a-zA-Z])\^([0-9]+)', r'\1^{\2}', latex_text) 
        latex_text = re.sub(r'([a-zA-Z])_\{([^}]+)\}', r'\1_{\2}', latex_text)
        latex_text = re.sub(r'([a-zA-Z])\^\{([^}]+)\}', r'\1^{\2}', latex_text)
        
        # Wrap in display math if it contains LaTeX commands
        if '\\' in latex_text and len(latex_text.strip()) > 2:
            return f"$${latex_text.strip()}$$"
        
        return None
    
    def _all_placeholders_replaced(self, text: str) -> bool:
        """Check if all visual placeholders have been replaced."""
        import re
        # Look for remaining placeholder patterns like [visual_0001]
        placeholder_pattern = r'`\[visual_\d+\]`'
        return not bool(re.search(placeholder_pattern, text))
    
    def _extract_page_specific_content(self, full_content: str, page_number: int) -> str:
        """Extract content specific to a particular page from the full VLM output."""
        import re
        
        # First, remove binary data (base64 encoded images)
        cleaned_content = self._remove_binary_data(full_content)
        
        # Split content by page markers or figure references
        if page_number == 1:
            # Look for content before "Figure 2" or page 2 markers
            patterns = [
                r'(.*?)(?=Figure 2:|## Page 2|\n2\n)',  # Content before Figure 2
                r'Figure 1[:\s]*(.*?)(?=Figure 2:|## Page 2|\n2\n)',  # Content around Figure 1
            ]
        elif page_number == 2:
            # Look for content after Figure 1 or page 1 content
            patterns = [
                r'Figure 2[:\s]*(.*?)(?=## Page 3|\n3\n|$)',  # Content around Figure 2
                r'(?:Figure 2:|## Page 2|\n2\n)(.*?)(?=## Page 3|\n3\n|$)',  # Content after page 2 marker
            ]
        else:
            return ""
            
        for pattern in patterns:
            match = re.search(pattern, cleaned_content, re.DOTALL | re.IGNORECASE)
            if match:
                content = match.group(1).strip()
                if len(content) > 50:  # Only return substantial content
                    return content
                    
        return ""
        
    def _remove_binary_data(self, content: str) -> str:
        """Remove binary data (base64 images, blobs) from content to keep only text descriptions."""
        import re
        
        # Remove base64 encoded images
        content = re.sub(r'data:image/[^;]+;base64,[A-Za-z0-9+/=]+', '[Image removed]', content)
        
        # Remove standalone base64 strings (long sequences of base64 characters)
        content = re.sub(r'[A-Za-z0-9+/]{100,}={0,2}', '[Binary data removed]', content)
        
        # Remove ![Image](...) markdown with data: URLs
        content = re.sub(r'!\[Image\]\(data:[^)]+\)', '[Image reference removed]', content)
        
        # Clean up multiple consecutive newlines
        content = re.sub(r'\n{3,}', '\n\n', content)
        
        return content
    
    def _remove_visual_elements_section(self, text: str) -> str:
        """Remove the Visual Elements section after integration is complete."""
        import re
        
        # Pattern to match the entire Visual Elements section
        pattern = r'\n---\n\n## Visual Elements\n.*?\*Enhanced content will be updated automatically when processing completes\.\*'
        enhanced_text = re.sub(pattern, '', text, flags=re.DOTALL)
        
        return enhanced_text

    def _extract_csv(self, file_path: Path) -> ExtractionResult:
        """Fast CSV text extraction."""
        try:
            # Try pandas first for complex CSVs
            df = pd.read_csv(file_path, nrows=1000)  # Limit rows for speed
            text_content = df.to_string()
            row_count = len(df)
            col_count = len(df.columns)
            
        except Exception:
            # Fallback to basic CSV reader
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                reader = csv.reader(f)
                rows = []
                for i, row in enumerate(reader):
                    if i > 1000:  # Limit for speed
                        break
                    rows.append(" | ".join(row))
                text_content = "\n".join(rows)
                row_count = len(rows)
                col_count = len(rows[0].split(',')) if rows else 0
        
        # CSV is inherently tabular
        placeholder = self._create_placeholder(
            "table",
            1,
            f"CSV data: {row_count} rows, {col_count} columns"
        )
        
        return ExtractionResult(
            success=True,
            text_content=text_content,
            metadata={
                "file_path": str(file_path),
                "row_count": row_count,
                "column_count": col_count
            },
            visual_placeholders=[placeholder],
            extraction_method=ExtractionMethod.BUILTIN,
            extraction_time=0.0,
            page_count=1,
            word_count=0,
            char_count=0
        )
    
    def _extract_text(self, file_path: Path) -> ExtractionResult:
        """Fast plain text extraction."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text_content = f.read()
        except Exception as e:
            return self._extraction_failed(f"Text extraction failed: {e}")
        
        return ExtractionResult(
            success=True,
            text_content=text_content,
            metadata={"file_path": str(file_path)},
            visual_placeholders=[],
            extraction_method=ExtractionMethod.BUILTIN,
            extraction_time=0.0,
            page_count=1,
            word_count=0,
            char_count=0
        )
    
    def _extract_json(self, file_path: Path) -> ExtractionResult:
        """Fast JSON text extraction."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            # Convert JSON to readable text
            text_content = json.dumps(data, indent=2, ensure_ascii=False)
            
        except Exception as e:
            return self._extraction_failed(f"JSON extraction failed: {e}")
        
        return ExtractionResult(
            success=True,
            text_content=text_content,
            metadata={"file_path": str(file_path)},
            visual_placeholders=[],
            extraction_method=ExtractionMethod.BUILTIN,
            extraction_time=0.0,
            page_count=1,
            word_count=0,
            char_count=0
        )
    
    def _extract_xml(self, file_path: Path) -> ExtractionResult:
        """Fast XML text extraction."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                xml_content = f.read()
            
            if BS4_AVAILABLE:
                soup = BeautifulSoup(xml_content, 'xml')
                text_content = soup.get_text()
            else:
                # Basic XML tag removal
                import re
                text_content = re.sub('<[^<]+?>', '', xml_content)
            
        except Exception as e:
            return self._extraction_failed(f"XML extraction failed: {e}")
        
        return ExtractionResult(
            success=True,
            text_content=text_content,
            metadata={"file_path": str(file_path)},
            visual_placeholders=[],
            extraction_method=ExtractionMethod.BEAUTIFULSOUP if BS4_AVAILABLE else ExtractionMethod.BUILTIN,
            extraction_time=0.0,
            page_count=1,
            word_count=0,
            char_count=0
        )
    
    def _embed_inline_placeholders(self, text: str, page_num: int, page_obj, images=None) -> tuple:
        """Embed visual placeholders inline where they actually appear in text."""
        import re
        
        enhanced_text = text
        placeholders = []
        
        # 1. Mathematical formulas - detect and replace with placeholders
        # Look for mathematical patterns (equations, symbols, etc.)
        math_patterns = [
            r'([∑∫∏√±∞≈≤≥≠∂∇∀∃∈∉⊂⊃∪∩∧∨¬→↔∴∵α-ωΑ-ΩΘΔΦΨΩπλμσφθγδε]+[^.!?\n]*)',
            r'(\b[xyz]_\{[^}]+\}|\b[xyz]\^\{[^}]+\})',  # Subscripts/superscripts
            r'(\$[^$\n]+\$|\\\[[^\]]+\\\])',  # LaTeX math
            r'([a-zA-Z]_[0-9]+|[a-zA-Z]\^[0-9]+)',  # Simple sub/superscripts
            r'((?:\b[ABCDEFGHIJKLMNOPQRSTUVWXYZ]_[a-zA-Z0-9]+\b[\s\+\-\=\(\)]*){2,})',  # Multi-variable equations
        ]
        
        for pattern in math_patterns:
            matches = list(re.finditer(pattern, enhanced_text))
            # Process in reverse order to avoid position shifts
            for match in reversed(matches):
                formula_text = match.group(1)
                if len(formula_text) > 3:  # Only process substantial formulas
                    placeholder = self._create_placeholder("formula", page_num, f"Mathematical formula: {formula_text[:50]}")
                    placeholder_tag = f"`[{placeholder['placeholder_id']}]`"
                    enhanced_text = enhanced_text[:match.start()] + placeholder_tag + enhanced_text[match.end():]
                    placeholders.append(placeholder)
        
        # 2. Tables - detect and replace with placeholders  
        table_patterns = [
            r'(\|[^\n]+\|[^\n]*\n\|[-\s:|]+\|[^\n]*\n(?:\|[^\n]+\|[^\n]*\n)*)',  # Markdown tables
            r'((?:^[^\n]*\t[^\n]*\n){3,})',  # Tab-separated tables (3+ rows)
            r'((?:[A-Za-z0-9\s]{2,}\s{2,}[A-Za-z0-9\s]{2,}\s{2,}[A-Za-z0-9\s]{2,}[\n\r]){2,})',  # Space-aligned tables
        ]
        
        for pattern in table_patterns:
            matches = list(re.finditer(pattern, enhanced_text, re.MULTILINE))
            for match in reversed(matches):
                table_content = match.group(1).strip()
                if len(table_content.split('\n')) >= 3:  # At least 3 rows
                    placeholder = self._create_placeholder("table", page_num, f"Table with {len(table_content.split())} entries")
                    placeholder_tag = f"`[{placeholder['placeholder_id']}]`"
                    enhanced_text = enhanced_text[:match.start()] + placeholder_tag + enhanced_text[match.end():]
                    placeholders.append(placeholder)
        
        # 3. Figure references - replace with image placeholders
        figure_patterns = [
            r'(Figure\s+(\d+):\s*([^.]+\.?))',
            r'(Fig\.?\s+(\d+)[\.:]\s*([^.]+\.?))',
            r'(Chart\s+(\d+)[\.:]\s*([^.]+\.?))',
            r'(Graph\s+(\d+)[\.:]\s*([^.]+\.?))',
        ]
        
        for pattern in figure_patterns:
            matches = list(re.finditer(pattern, enhanced_text, re.IGNORECASE))
            for match in reversed(matches):
                full_caption = match.group(1)
                fig_num = match.group(2) if len(match.groups()) > 1 else "1"
                description = match.group(3) if len(match.groups()) > 2 else "Figure"
                
                # Determine element type from description
                element_type = "image"
                desc_lower = description.lower()
                if any(word in desc_lower for word in ["chart", "graph", "bar", "plot", "histogram"]):
                    element_type = "chart"
                elif any(word in desc_lower for word in ["diagram", "flow", "schematic"]):
                    element_type = "diagram"
                elif any(word in desc_lower for word in ["map", "location", "geographic"]):
                    element_type = "diagram"
                
                placeholder = self._create_placeholder(element_type, page_num, full_caption)
                placeholder_tag = f"`[{placeholder['placeholder_id']}]`"
                enhanced_text = enhanced_text[:match.start()] + placeholder_tag + enhanced_text[match.end():]
                placeholders.append(placeholder)
        
        # 4. Algorithm blocks - replace with special placeholders
        algorithm_pattern = r'(Algorithm\s+\d+[^\n]*\n(?:[^\n]+\n)*?(?:return|end|until).*)'
        matches = list(re.finditer(algorithm_pattern, enhanced_text, re.IGNORECASE | re.MULTILINE))
        for match in reversed(matches):
            algorithm_text = match.group(1)
            placeholder = self._create_placeholder("diagram", page_num, f"Algorithm block: {algorithm_text[:50]}...")
            placeholder_tag = f"`[{placeholder['placeholder_id']}]`"
            enhanced_text = enhanced_text[:match.start()] + placeholder_tag + enhanced_text[match.end():]
            placeholders.append(placeholder)
        
        # 5. If we have images but no figure references, add placeholders for orphaned images
        if images and not any(p['element_type'] in ['image', 'chart', 'diagram'] for p in placeholders):
            for i, img in enumerate(images):
                placeholder = self._create_placeholder("image", page_num, f"Image {i+1} on page {page_num}")
                placeholders.append(placeholder)
                # Don't embed these inline since we don't know where they belong
        
        return enhanced_text, placeholders

    def _detect_special_regions(self, page, page_num: int) -> tuple:
        """Detect mathematical, table, and figure regions before text extraction."""
        import re
        
        math_regions = []
        table_regions = []
        figure_regions = []
        
        # STEP 1: Analyze text blocks for mathematical and table content
        text_dict = page.get_text("dict")
        
        for block in text_dict.get("blocks", []):
            if "lines" not in block:
                continue
                
            block_text = ""
            for line in block["lines"]:
                for span in line.get("spans", []):
                    block_text += span.get("text", "")
                block_text += "\n"
            
            block_bbox = block["bbox"]  # [x0, y0, x1, y1]
            
            # 1. Detect mathematical expressions by font and symbols
            if self._is_mathematical_content(block, block_text):
                placeholder = self._create_placeholder("formula", page_num, f"Mathematical expression: {block_text[:50]}")
                placeholder["bbox"] = [float(block_bbox[0]), float(block_bbox[1]), float(block_bbox[2]), float(block_bbox[3])]
                placeholder["original_text"] = block_text
                
                # EXTRACT MATHEMATICAL REGION AS IMAGE
                self._extract_bbox_as_image(page, block_bbox, placeholder, page_num, "formula")
                
                math_regions.append(placeholder)
            
            # 2. Detect table-like structures
            elif self._is_table_content(block_text, block):
                placeholder = self._create_placeholder("table", page_num, f"Table structure")
                placeholder["bbox"] = [float(block_bbox[0]), float(block_bbox[1]), float(block_bbox[2]), float(block_bbox[3])]
                placeholder["original_text"] = block_text
                
                # EXTRACT TABLE REGION AS IMAGE
                self._extract_bbox_as_image(page, block_bbox, placeholder, page_num, "table")
                
                table_regions.append(placeholder)
        
        # STEP 2: DISABLED - Skip embedded images/charts for now
        # We're focusing only on formulas and tables
        self._log(f"🚫 Skipping embedded image detection (focusing on formulas/tables only)")
        
        # STEP 3: DISABLED - Skip figure captions for now
        # We're focusing only on formulas and tables  
        self._log(f"🚫 Skipping figure caption detection (focusing on formulas/tables only)")
        
        # STEP 4: DISABLED - Keep ALL detected regions (even overlapping)
        # This ensures we catch everything, even if some regions overlap
        self._log(f"📋 Keeping all detected regions without merging:")
        
        return math_regions, table_regions, figure_regions
    
    def _classify_image_by_context(self, page, img_bbox, text_dict) -> str:
        """Classify image type based on nearby text context."""
        import re
        
        # Look for text blocks near the image
        nearby_text = []
        
        for block in text_dict.get("blocks", []):
            if "lines" not in block:
                continue
            
            block_bbox = block["bbox"]
            
            # Check if text block is near the image (within reasonable distance)
            if self._is_near_bbox(img_bbox, block_bbox, threshold=100):  # 100 pixel threshold
                block_text = ""
                for line in block["lines"]:
                    for span in line.get("spans", []):
                        block_text += span.get("text", "")
                    block_text += " "
                nearby_text.append(block_text.strip())
        
        # Analyze nearby text to classify image type
        combined_text = " ".join(nearby_text).lower()
        
        # Chart/Graph indicators
        if any(keyword in combined_text for keyword in 
               ['chart', 'graph', 'plot', 'axis', 'data', 'trend', 'distribution', 
                'percentage', '%', 'bar', 'pie', 'line chart', 'scatter']):
            return "chart"
        
        # Diagram indicators
        elif any(keyword in combined_text for keyword in 
                 ['diagram', 'flowchart', 'schema', 'architecture', 'workflow', 
                  'process', 'system', 'network', 'structure']):
            return "diagram"
        
        # Algorithm/Code indicators
        elif any(keyword in combined_text for keyword in 
                 ['algorithm', 'pseudocode', 'code', 'function', 'procedure']):
            return "algorithm"
        
        # Default to generic figure
        else:
            return "figure"
    
    def _is_near_bbox(self, bbox1, bbox2, threshold=50) -> bool:
        """Check if two bounding boxes are near each other."""
        # Calculate center points
        center1_x = (bbox1[0] + bbox1[2]) / 2
        center1_y = (bbox1[1] + bbox1[3]) / 2
        center2_x = (bbox2[0] + bbox2[2]) / 2
        center2_y = (bbox2[1] + bbox2[3]) / 2
        
        # Calculate distance
        distance = ((center1_x - center2_x) ** 2 + (center1_y - center2_y) ** 2) ** 0.5
        
        return distance <= threshold
    
    def _is_mathematical_content(self, block, text: str) -> bool:
        """Detect if a text block contains mathematical content."""
        import re
        
        # Check for mathematical symbols (Unicode)
        math_symbols = ['∑', '∫', '∏', '√', '±', '∞', '≈', '≤', '≥', '≠', '∂', '∇', 
                       '∀', '∃', '∈', '∉', '⊂', '⊃', '∪', '∩', '∧', '∨', '¬', '→', '↔',
                       'α', 'β', 'γ', 'δ', 'ε', 'ζ', 'η', 'θ', 'λ', 'μ', 'π', 'σ', 'φ', 'ψ', 'ω']
        
        if any(symbol in text for symbol in math_symbols):
            return True
        
        # Check for common math patterns (VERY AGGRESSIVE DETECTION)
        math_patterns = [
            # Basic equations and expressions
            r'\b[a-zA-Z]\s*[=≈]\s*[0-9\-+\s\(\)]+',  # x = 123, y = -5, etc.
            r'\([^)]*[+\-*/][^)]*\)',  # (a + b), (x - y), etc.
            r'\b[a-zA-Z]\^[0-9]+\b',  # x^2, a^3
            r'\b[a-zA-Z]_[0-9]+\b',   # x_1, y_2  
            r'\b\d+\s*[+\-×÷*/]\s*\d+',  # 5 + 3, 2 * 4
            
            # Mathematical functions
            r'\bsin\b|\bcos\b|\btan\b|\blog\b|\bln\b|\bexp\b',  # trig/log functions
            r'\b[f|g|h]\([^)]*\)',    # function notation f(x)
            r'\bmax\b|\bmin\b|\bsum\b|\bprod\b',  # mathematical operators
            
            # Advanced math
            r'\bd[xy]/d[xy]\b',       # derivatives
            r'∂[^∂]*∂',               # partial derivatives
            r'\bE\[.*\]',            # expectation notation
            r'\bP\(.*\)',            # probability notation
            
            # AGGRESSIVE: Any line with variables and operators
            r'\b[a-zA-Z]\s*[+\-*/=]\s*[a-zA-Z0-9]',  # a + b, x = y, etc.
            r'\b\d+[a-zA-Z]',        # 2x, 3y (coefficients)
            r'\b[a-zA-Z]\d+',        # x2, y3 (variables with numbers)
            r'[≤≥<>]\s*\d+',         # inequalities
            r'\b\d+\.\d+\b',         # decimal numbers (often in formulas)
            r'\b[A-Z][a-z]*\s*=',    # Variable definitions (Velocity =, etc.)
            r'\(\d+\)',              # numbered equations (1), (2)
            r'\[[0-9]+\]',           # equation references [1], [2]
            
            # Mathematical relationships
            r'\bif\s+.*\bthen\b',    # conditional statements
            r'\bwhere\s+[a-zA-Z]',   # variable definitions
            r'\blet\s+[a-zA-Z]',     # variable declarations
            r'\bfor\s+[a-zA-Z]\s*[∈=]',  # for loops in math
            r'\∀|\∃',                # quantifiers
        ]
        
        for pattern in math_patterns:
            if re.search(pattern, text, re.IGNORECASE):
                return True
        
        # Check for subscript/superscript patterns
        if re.search(r'[a-zA-Z][₀₁₂₃₄₅₆₇₈₉]|[a-zA-Z][⁰¹²³⁴⁵⁶⁷⁸⁹]', text):
            return True
            
        # Check for font characteristics (italic math fonts)
        math_font_count = 0
        total_spans = 0
        
        for line in block.get("lines", []):
            for span in line.get("spans", []):
                total_spans += 1
                font = span.get("font", "").lower()
                # Common math fonts in PDFs
                if any(math_font in font for math_font in ["math", "italic", "symbol", "times-italic", "cambria-math"]):
                    math_font_count += 1
        
        # If >20% of text uses math fonts, likely mathematical (lowered threshold)
        if total_spans > 0 and math_font_count / total_spans > 0.2:
            return True
        
        return False
    
    def _is_table_content(self, text: str, block) -> bool:
        """Detect if content represents a table structure."""
        import re
        
        lines = text.strip().split('\n')
        
        # Basic length check
        if len(lines) < 2:
            return False
        
        # Check for explicit table indicators
        table_keywords = ['table', 'column', 'row', 'header', 'data', 'results']
        if any(keyword in text.lower() for keyword in table_keywords):
            # If table keywords found, check for structured data
            if self._has_structured_layout(lines, block):
                return True
        
        # Look for consistent column structures
        if len(lines) >= 3:
            # Check for markdown-style tables
            if any('|' in line for line in lines):
                pipe_counts = [line.count('|') for line in lines if '|' in line]
                if len(set(pipe_counts)) == 1 and pipe_counts[0] >= 2:  # Consistent pipe count
                    return True
            
            # Check for space-aligned columns (at least 2 columns)
            tab_counts = [line.count('\t') for line in lines]
            if any(count >= 1 for count in tab_counts):  # Lowered threshold
                return True
            
            # Check for consistent spacing patterns (table-like alignment)
            spacing_patterns = []
            for line in lines:
                # Find sequences of multiple spaces (potential column separators)
                spaces = re.findall(r'\s{2,}', line)
                spacing_patterns.append(len(spaces))
            
            # If most lines have similar spacing patterns, likely a table
            if spacing_patterns and len(set(spacing_patterns)) <= 2 and max(spacing_patterns) >= 2:
                return True
            
            # Check for numeric data in columns
            numeric_lines = 0
            for line in lines:
                numbers_in_line = len(re.findall(r'\b\d+\.?\d*\b', line))
                if numbers_in_line >= 2:  # At least 2 numbers per line
                    numeric_lines += 1
            
            # If >30% of lines contain multiple numbers, likely tabular data (lowered threshold)
            if numeric_lines / len(lines) > 0.3:
                return True
            
        # AGGRESSIVE: Any text block with consistent indentation might be tabular
        if len(lines) >= 2:
            indentations = []
            for line in lines:
                leading_spaces = len(line) - len(line.lstrip())
                indentations.append(leading_spaces)
            
            # If we have 2+ different consistent indentation levels, might be structured data
            unique_indents = len(set(indentations))
            if unique_indents >= 2 and unique_indents <= 4:  # 2-4 indentation levels
                return True
        
        return False
    
    def _has_structured_layout(self, lines: List[str], block) -> bool:
        """Check if text block has structured layout suggesting a table."""
        # Look at font/positioning consistency across lines
        if not block.get("lines"):
            return False
        
        # Check for consistent left margins (aligned columns)
        left_margins = []
        for line_data in block["lines"]:
            if line_data.get("spans"):
                first_span = line_data["spans"][0]
                left_margins.append(first_span.get("bbox", [0])[0])
        
        # If we have consistent left alignment, check for other column positions
        if len(set(left_margins)) <= 2:  # Maximum 2 different left margins
            return True
            
        return False
    
    def _is_figure_content(self, text: str) -> bool:
        """Detect figure captions and references."""
        import re
        
        # Look for figure/chart/algorithm patterns
        patterns = [
            r'Figure\s+\d+[:.]\s*',
            r'Fig\.?\s+\d+[:.]\s*',
            r'Chart\s+\d+[:.]\s*',
            r'Algorithm\s+\d+\s*',
            r'Table\s+\d+[:.]\s*',
        ]
        
        return any(re.search(pattern, text, re.IGNORECASE) for pattern in patterns)
    
    def _determine_figure_type(self, text: str) -> str:
        """Determine the specific type of figure based on content analysis."""
        import re
        
        text_lower = text.lower()
        
        # Algorithm blocks - highest priority
        if re.search(r'algorithm\s+\d+', text, re.IGNORECASE):
            return "diagram"
        
        # Charts and graphs - look for data visualization keywords
        chart_keywords = [
            'chart', 'graph', 'plot', 'histogram', 'bar chart', 'pie chart', 
            'scatter plot', 'line graph', 'data', 'scores', 'results',
            'performance', 'comparison', 'metrics', 'rate', 'distribution'
        ]
        
        if any(keyword in text_lower for keyword in chart_keywords):
            return "chart"
        
        # Diagrams - look for structural/flow keywords  
        diagram_keywords = [
            'diagram', 'model', 'architecture', 'flow', 'process', 'structure',
            'graphical model', 'network', 'framework', 'schematic'
        ]
        
        if any(keyword in text_lower for keyword in diagram_keywords):
            return "diagram"
        
        # Sample images - look for image/sample keywords
        image_keywords = [
            'samples', 'images', 'generated', 'interpolation', 'reconstruction',
            'faces', 'pictures', 'visual', 'example'
        ]
        
        if any(keyword in text_lower for keyword in image_keywords):
            return "image"
        
        # Default to image for generic figures
        return "image"
    
    def _extract_text_with_preservation(self, page, math_regions, table_regions, figure_regions, page_num: int) -> str:
        """Extract text while replacing only mathematical expressions inline, preserving surrounding text."""
        
        # Get all text from the page
        full_text = page.get_text()
        
        # Handle different types of regions differently
        
        # 1. Replace complete figures and tables (these should be standalone)
        standalone_regions = table_regions + [r for r in figure_regions if r.get('element_type') in ['image', 'chart', 'diagram']]
        standalone_regions.sort(key=lambda r: r["bbox"][1], reverse=True)  # Bottom to top
        
        for region in standalone_regions:
            region_text = region["original_text"].strip()
            if region_text in full_text:
                placeholder_tag = f"`[{region['placeholder_id']}]`"
                full_text = full_text.replace(region_text, placeholder_tag, 1)
        
        # 2. Handle mathematical formulas inline - preserve surrounding text
        math_regions.sort(key=lambda r: r["bbox"][1], reverse=True)  # Bottom to top
        
        for region in math_regions:
            full_text = self._replace_math_inline(full_text, region)
        
        return full_text
    
    def _replace_math_inline(self, text: str, math_region) -> str:
        """Replace mathematical expressions inline while preserving surrounding text."""
        import re
        
        original_text = math_region["original_text"].strip()
        placeholder_id = math_region["placeholder_id"]
        
        # If the math region is a whole paragraph, we need to extract just the mathematical parts
        if len(original_text) > 100:  # Likely a paragraph with embedded math
            # Find mathematical expressions within the text
            math_expressions = self._find_math_expressions_in_text(original_text)
            
            if math_expressions:
                # Replace each mathematical expression with a placeholder
                modified_text = original_text
                for i, math_expr in enumerate(reversed(math_expressions)):  # Reverse to preserve positions
                    math_placeholder = f"`[{placeholder_id}_part{len(math_expressions)-i}]`"
                    modified_text = modified_text.replace(math_expr, math_placeholder, 1)
                
                # Replace the original paragraph with the modified version
                if original_text in text:
                    text = text.replace(original_text, modified_text, 1)
                else:
                    # Fallback: try to find a close match
                    # Look for the first 50 and last 50 characters
                    start_chunk = original_text[:50]
                    end_chunk = original_text[-50:]
                    
                    # Find the paragraph in the full text
                    start_pos = text.find(start_chunk)
                    if start_pos >= 0:
                        end_pos = text.find(end_chunk, start_pos)
                        if end_pos >= 0:
                            end_pos += len(end_chunk)
                            # Replace the found paragraph
                            text = text[:start_pos] + modified_text + text[end_pos:]
            
        else:
            # Small math expression - replace entirely
            if original_text in text:
                placeholder_tag = f"`[{placeholder_id}]`"
                text = text.replace(original_text, placeholder_tag, 1)
        
        return text
    
    def _find_math_expressions_in_text(self, text: str) -> list:
        """Find individual mathematical expressions within a text block."""
        import re
        
        math_expressions = []
        
        # Look for mathematical symbols and expressions
        patterns = [
            r'[∑∫∏√±∞≈≤≥≠∂∇∀∃∈∉⊂⊃∪∩∧∨¬→↔∴∵α-ωΑ-ΩΘΔΦΨΩπλμσφθγδε]+[^\s.!?]*',  # Unicode math symbols
            r'[a-zA-Z][₀₁₂₃₄₅₆₇₈₉]+[^\s.!?]*',  # Subscripts
            r'[a-zA-Z][⁰¹²³⁴⁵⁶⁷⁸⁹]+[^\s.!?]*',  # Superscripts
            r'ρ\$\\_{[^}]+}\$\s*\([^)]+\)',  # Specific notation like ρ${θ}$ ( x${0}$ )
            r'\$[^$]+\$',  # LaTeX-style math
            r'[A-Z]\s*\([^)]+\)',  # Function notation like N ( x ; 0 ; I )
            r'[x-z]\\?_\{[^}]+\}',  # Variable subscripts
            r':=\s*∫[^,]+',  # Integral definitions
        ]
        
        for pattern in patterns:
            matches = re.finditer(pattern, text)
            for match in matches:
                expr = match.group().strip()
                if len(expr) > 2 and expr not in math_expressions:
                    math_expressions.append(expr)
        
        # Sort by length (longest first) to avoid replacing substrings
        math_expressions.sort(key=len, reverse=True)
        
        return math_expressions
    
    def _merge_overlapping_regions(self, regions: List[Dict], page_num: int, region_type: str) -> List[Dict]:
        """Merge overlapping regions to prevent duplicate extractions."""
        if not regions:
            return regions
        
        # Sort regions by vertical position (top to bottom)
        regions.sort(key=lambda r: r["bbox"][1])  # Sort by y0 coordinate
        
        merged_regions = []
        current_group = [regions[0]]
        
        for region in regions[1:]:
            current_bottom = current_group[-1]["bbox"][3]  # y1 of last region in group
            region_top = region["bbox"][1]  # y0 of current region
            
            # Check for vertical overlap or proximity (within 40 points)
            if region_top <= current_bottom + 40:
                current_group.append(region)
            else:
                # Process current group
                if len(current_group) > 1:
                    merged_region = self._create_merged_region(current_group, page_num, region_type)
                    merged_regions.append(merged_region)
                else:
                    merged_regions.extend(current_group)
                
                # Start new group
                current_group = [region]
        
        # Process final group
        if len(current_group) > 1:
            merged_region = self._create_merged_region(current_group, page_num, region_type)
            merged_regions.append(merged_region)
        else:
            merged_regions.extend(current_group)
        
        print(f"📋 {region_type.title()} regions: {len(regions)} → {len(merged_regions)} (after merging)")
        return merged_regions
    
    def _create_merged_region(self, regions: List[Dict], page_num: int, region_type: str) -> Dict:
        """Create a single merged region from multiple overlapping regions."""
        # Calculate merged bounding box
        min_x = min(r["bbox"][0] for r in regions)
        min_y = min(r["bbox"][1] for r in regions)
        max_x = max(r["bbox"][2] for r in regions)
        max_y = max(r["bbox"][3] for r in regions)
        
        # Combine original text
        combined_text = "\n".join(r.get("original_text", "") for r in regions if r.get("original_text"))
        
        # Create merged placeholder
        merged_region = self._create_placeholder(
            region_type, 
            page_num, 
            f"Merged {region_type} region ({len(regions)} elements)"
        )
        
        merged_region["bbox"] = [float(min_x), float(min_y), float(max_x), float(max_y)]
        merged_region["original_text"] = combined_text
        merged_region["merged_from"] = [r["placeholder_id"] for r in regions]
        
        return merged_region
    
    def _extract_bbox_as_image(self, page, bbox, placeholder, page_num: int, element_type: str):
        """Extract the bounding box region as an image file and save to output folder."""
        # Skip image extraction in text-only mode
        if self.text_only_mode:
            self._log(f"⚡ Text-only mode: skipping image extraction for {element_type}")
            return
        
        try:
            # Create output directory
            output_dir = Path('/home/corey/projects/docling/cli/output/latest')
            output_dir.mkdir(parents=True, exist_ok=True)
            
            # Generate filename
            placeholder_id = placeholder["placeholder_id"]
            filename = f"{placeholder_id}_page{page_num}_{element_type}.png"
            output_path = output_dir / filename
            
            # NORMALIZE TO FULL PAGE WIDTH WITH PADDING
            page_rect = page.rect
            page_width = page_rect.width
            page_height = page_rect.height
            
            # Convert bbox from list/tuple [x0, y0, x1, y1] to Rect if needed
            if isinstance(bbox, (list, tuple)):
                bbox_rect = fitz.Rect(bbox[0], bbox[1], bbox[2], bbox[3])
            else:
                bbox_rect = bbox
            
            # Expand to full page width, add vertical padding for context
            normalized_bbox = fitz.Rect(
                0,  # Start at left margin
                max(0, bbox_rect.y0 - 20),  # Add 20pt padding above
                page_width,  # Extend to full page width
                min(page_height, bbox_rect.y1 + 20)  # Add 20pt padding below
            )
            
            # Extract the normalized region as pixmap (image)
            mat = fitz.Matrix(2, 2)  # 2x scaling for better quality
            pix = page.get_pixmap(matrix=mat, clip=normalized_bbox)
            
            # Save as PNG
            pix.save(str(output_path))
            pix = None  # Free memory
            
            # Store the extracted image path and normalized bbox in placeholder
            placeholder["extracted_image_path"] = str(output_path)
            placeholder["normalized_bbox"] = [float(normalized_bbox.x0), float(normalized_bbox.y0), 
                                            float(normalized_bbox.x1), float(normalized_bbox.y1)]
            
            original_height = bbox_rect.y1 - bbox_rect.y0
            normalized_height = normalized_bbox.y1 - normalized_bbox.y0
            
            self._log(f"💾 Extracted {element_type} image: {filename}")
            self._log(f"   📏 Original: {bbox_rect.width:.0f}x{original_height:.0f} → Normalized: {page_width:.0f}x{normalized_height:.0f}")
            
        except Exception as e:
            self._log(f"⚠️  Failed to extract image for {placeholder_id}: {e}")
            placeholder["extraction_error"] = str(e)

    def _create_placeholder(self, element_type: str, page_number: int, description: str) -> Dict[str, Any]:
        """Create a visual placeholder for later processing."""
        self.placeholder_counter += 1
        
        return {
            "placeholder_id": f"visual_{self.placeholder_counter:04d}",
            "element_type": element_type,
            "page_number": page_number,
            "description": description,
            "priority": 3,
            "status": "queued"
        }
    
    def _extraction_failed(self, error_message: str) -> ExtractionResult:
        """Helper for failed extraction results."""
        return ExtractionResult(
            success=False,
            text_content="",
            metadata={},
            visual_placeholders=[],
            extraction_method=ExtractionMethod.BUILTIN,
            extraction_time=0.0,
            page_count=0,
            word_count=0,
            char_count=0,
            error_message=error_message
        )
    
    def generate_markdown_with_placeholders(self, result: ExtractionResult, 
                                           file_path: Path) -> str:
        """
        Generate markdown output with visual placeholders.
        
        This provides immediate results while visual processing is queued.
        """
        
        markdown_lines = []
        
        # Add title
        title = file_path.stem.replace('_', ' ').replace('-', ' ').title()
        markdown_lines.append(f"# {title}")
        markdown_lines.append("")
        
        # Add extraction metadata
        markdown_lines.append("<!-- Docling Fast Extraction -->")
        markdown_lines.append(f"<!-- Extracted: {time.strftime('%Y-%m-%d %H:%M:%S')} -->")
        markdown_lines.append(f"<!-- Method: {result.extraction_method.value} -->")
        markdown_lines.append(f"<!-- Time: {result.extraction_time:.3f}s -->")
        markdown_lines.append("")
        
        # Add visual processing status if applicable
        if result.visual_placeholders:
            markdown_lines.append("<!-- Visual Enhancement Status -->")
            markdown_lines.append(f"<!-- {len(result.visual_placeholders)} visual elements queued for processing -->")
            markdown_lines.append("")
        
        # Add main content
        if result.text_content:
            markdown_lines.append(result.text_content)
        else:
            markdown_lines.append("*No text content extracted*")
        
        # Add visual placeholders section for pending processing
        if result.visual_placeholders:
            markdown_lines.append("")
            markdown_lines.append("---")
            markdown_lines.append("")
            markdown_lines.append("## Visual Elements")
            markdown_lines.append("")
            markdown_lines.append("*The following visual elements are being processed for enhancement:*")
            markdown_lines.append("")
            
            for placeholder in result.visual_placeholders:
                element_name = placeholder['element_type'].title()
                page_num = placeholder['page_number']
                description = placeholder.get('description', 'Processing...')
                placeholder_id = placeholder['placeholder_id']
                markdown_lines.append(f"- **{element_name}** (Page {page_num}): {description} `[{placeholder_id}]`")
            
            markdown_lines.append("")
            markdown_lines.append("*Enhanced content will be updated automatically when processing completes.*")
        
        return "\n".join(markdown_lines)


def main():
    """Extract text from specified file or run test on data directory."""
    import sys
    
    extractor = FastTextExtractor(log_to_file=True)
    
    # Check if specific file was provided
    if len(sys.argv) > 1:
        # Process specific file
        file_path = Path(sys.argv[1])
        if not file_path.exists():
            print(f"❌ File not found: {file_path}")
            return
        
        print("⚡ FAST TEXT EXTRACTION")
        print("=" * 50)
        print(f"\n📄 Extracting: {file_path.name}")
        
        result = extractor.extract(file_path)
        
        if result.success:
            pages_per_sec = result.page_count / result.extraction_time if result.extraction_time > 0 else float('inf')
            print(f"   ✅ Success: {result.extraction_method.value}")
            print(f"   📊 Pages: {result.page_count}")
            print(f"   📝 Words: {result.word_count:,}")
            print(f"   📄 Characters: {result.char_count:,}")
            print(f"   ⚡ Speed: {pages_per_sec:.1f} pages/sec")
            print(f"   ⏱️  Time: {result.extraction_time:.3f}s")
            print(f"   🎨 Visual placeholders: {len(result.visual_placeholders)}")
            
            if result.visual_placeholders:
                print(f"\n🎨 Visual Elements Found:")
                for i, placeholder in enumerate(result.visual_placeholders[:5], 1):
                    print(f"   {i}. {placeholder}")
            
            # Show first 500 characters of extracted text
            preview = result.text_content[:500]
            if len(result.text_content) > 500:
                preview += "..."
            print(f"\n📝 Text Preview:\n{preview}")
            
            # Show log file location
            if extractor.log_file:
                print(f"\n📋 Detailed log saved to: {extractor.log_file}")
                print(f"💾 Extracted images saved to: /home/corey/projects/docling/cli/output/latest/")
            
        else:
            print(f"   ❌ Failed: {result.error_message}")
        
        return
    
    # Default: Run test on sample files
    data_dir = Path('/home/corey/projects/docling/cli/data')
    test_files = []
    
    # Collect various file types
    for pattern in ['*.pdf', '*.docx', '*.html', '*.md', '*.csv']:
        test_files.extend(list(data_dir.glob(f'**/{pattern}'))[:2])  # 2 of each type
    
    if not test_files:
        print("No test files found")
        return
    
    print("⚡ FAST TEXT EXTRACTION TEST")
    print("=" * 50)
    
    total_start = time.time()
    
    for file_path in test_files:
        print(f"\n📄 Extracting: {file_path.name}")
        
        result = extractor.extract(file_path)
        
        if result.success:
            pages_per_sec = result.page_count / result.extraction_time if result.extraction_time > 0 else float('inf')
            print(f"   ✅ Success: {result.extraction_method.value}")
            print(f"   📊 Pages: {result.page_count}")
            print(f"   📝 Words: {result.word_count:,}")
            print(f"   ⚡ Speed: {pages_per_sec:.1f} pages/sec")
            print(f"   ⏱️  Time: {result.extraction_time:.3f}s")
            print(f"   🎨 Visual elements: {len(result.visual_placeholders)}")
        else:
            print(f"   ❌ Failed: {result.error_message}")
    
    total_time = time.time() - total_start
    print(f"\n📊 SUMMARY:")
    print(f"   Files processed: {len(test_files)}")
    print(f"   Total time: {total_time:.2f}s")
    print(f"   Average: {len(test_files)/total_time:.1f} files/sec")


if __name__ == "__main__":
    main()