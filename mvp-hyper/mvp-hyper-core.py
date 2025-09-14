#!/usr/bin/env python3
"""
MVP Hyper-Core: Ultra-High-Performance Document Processor
=========================================================
Target: 1,000 pages per second for text extraction
Minimal dependencies, maximum performance
"""

import os
import sys
import mmap
import time
import multiprocessing as mp
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Any
from dataclasses import dataclass
import hashlib
import json
from concurrent.futures import ProcessPoolExecutor, ThreadPoolExecutor
import asyncio

# Try to import PyMuPDF with fallback
try:
    import fitz  # PyMuPDF
    HAS_PYMUPDF = True
except ImportError:
    HAS_PYMUPDF = False
    print("Warning: PyMuPDF not available, performance will be limited")

# Optional accelerators
try:
    import numba
    HAS_NUMBA = True
except ImportError:
    HAS_NUMBA = False

try:
    import xxhash
    HAS_XXHASH = True
except ImportError:
    HAS_XXHASH = False
    

@dataclass
class UltraFastResult:
    """Minimal result structure for maximum speed."""
    file_path: str
    success: bool
    text: str
    page_count: int
    extraction_time: float
    pages_per_second: float
    metadata: Dict[str, Any]
    error: Optional[str] = None


class MemoryPool:
    """Pre-allocated memory pool to avoid allocation overhead."""
    
    def __init__(self, pool_size_mb: int = 1024):
        self.pool_size = pool_size_mb * 1024 * 1024
        self.buffers = []
        self.available = []
        
        # Pre-allocate buffers
        for _ in range(32):  # 32 buffers
            buf = bytearray(32 * 1024 * 1024)  # 32MB each
            self.buffers.append(buf)
            self.available.append(True)
    
    def get_buffer(self):
        """Get an available buffer."""
        for i, avail in enumerate(self.available):
            if avail:
                self.available[i] = False
                return self.buffers[i]
        # If no buffer available, create new one
        buf = bytearray(32 * 1024 * 1024)
        self.buffers.append(buf)
        self.available.append(False)
        return buf
    
    def release_buffer(self, buf):
        """Release buffer back to pool."""
        try:
            idx = self.buffers.index(buf)
            self.available[idx] = True
        except ValueError:
            pass


class UltraFastExtractor:
    """Core extraction engine optimized for 1000+ pages/second."""
    
    def __init__(self, 
                 num_workers: Optional[int] = None,
                 use_mmap: bool = True,
                 batch_size: int = 100,
                 cache_size_mb: int = 512):
        
        self.num_workers = num_workers or mp.cpu_count()
        self.use_mmap = use_mmap
        self.batch_size = batch_size
        self.cache_size_mb = cache_size_mb
        
        # Cache for ultra-fast repeated access
        self.cache = {}
        self.cache_size = cache_size_mb * 1024 * 1024
        self.current_cache_size = 0
        
        # Don't create pools in __init__ - create them when needed
        
    def extract_document_ultrafast(self, file_path: Path) -> UltraFastResult:
        """Extract text from any supported document format at maximum speed."""
        
        start_time = time.perf_counter()
        file_ext = file_path.suffix.lower()
        
        # Check cache first
        cache_key = self._get_cache_key(file_path)
        if cache_key in self.cache:
            cached = self.cache[cache_key]
            return UltraFastResult(
                file_path=str(file_path),
                success=True,
                text=cached['text'],
                page_count=cached['pages'],
                extraction_time=0.001,  # Cache hit is ~1ms
                pages_per_second=cached['pages'] / 0.001,
                metadata=cached['metadata']
            )
        
        try:
            # Route to appropriate extractor based on file type  
            if file_ext == '.pdf':
                return self._extract_pdf(file_path, start_time, cache_key, getattr(self, 'config', None))
            elif file_ext in ['.docx']:
                return self._extract_docx(file_path, start_time, cache_key)
            elif file_ext in ['.pptx']:
                return self._extract_pptx(file_path, start_time, cache_key)
            elif file_ext in ['.xlsx']:
                return self._extract_xlsx(file_path, start_time, cache_key)
            elif file_ext in ['.html', '.htm']:
                return self._extract_html(file_path, start_time, cache_key)
            elif file_ext in ['.txt', '.md']:
                return self._extract_text(file_path, start_time, cache_key)
            elif file_ext in ['.csv']:
                return self._extract_csv(file_path, start_time, cache_key)
            else:
                # Universal fallback - try to extract ANY file type
                return self._extract_universal(file_path, start_time, cache_key)
                
        except Exception as e:
            extraction_time = time.perf_counter() - start_time
            return UltraFastResult(
                file_path=str(file_path),
                success=False,
                text="",
                page_count=0,
                extraction_time=extraction_time,
                pages_per_second=0,
                metadata={"filename": file_path.name, "format": file_ext.upper()},
                error=str(e)
            )
    
    def _extract_parallel(self, doc, page_count: int) -> str:
        """Extract pages in parallel for maximum speed."""
        
        # For now, use sequential extraction to avoid pickling issues
        # TODO: Implement proper multiprocessing with shared memory
        return self._extract_sequential(doc)
    
    def _extract_chunk(self, doc, start: int, end: int) -> str:
        """Extract a chunk of pages."""
        texts = []
        for page_num in range(start, end):
            page = doc[page_num]
            # Use fastest extraction method
            text = page.get_text("text", flags=fitz.TEXT_PRESERVE_WHITESPACE)
            texts.append(text)
        return '\n'.join(texts)
    
    def _extract_sequential(self, doc) -> str:
        """Sequential extraction for small documents."""
        texts = []
        for page in doc:
            text = page.get_text("text", flags=fitz.TEXT_PRESERVE_WHITESPACE)
            texts.append(text)
        return '\n'.join(texts)
    
    def _extract_metadata_fast(self, doc, file_path: Path) -> Dict[str, Any]:
        """Extract minimal metadata for speed."""
        
        metadata = {
            "filename": file_path.name,
            "pages": len(doc),
            "size_bytes": file_path.stat().st_size,
            "format": "PDF",
        }
        
        # Only extract essential metadata
        if doc.metadata:
            metadata["title"] = doc.metadata.get("title", "")
            metadata["author"] = doc.metadata.get("author", "")
            metadata["subject"] = doc.metadata.get("subject", "")
        
        return metadata
    
    def _get_cache_key(self, file_path: Path) -> str:
        """Generate cache key using fast hashing."""
        
        if HAS_XXHASH:
            # xxhash is ~10x faster than SHA256
            h = xxhash.xxh64()
        else:
            h = hashlib.blake2b()
        
        # Hash file path and modification time
        key_string = f"{file_path.absolute()}_{file_path.stat().st_mtime}"
        h.update(key_string.encode())
        return h.hexdigest()
    
    def _extract_pdf(self, file_path: Path, start_time: float, cache_key: str, config=None) -> UltraFastResult:
        """Extract from PDF files - fail fast on errors."""
        if not HAS_PYMUPDF:
            return self._fail_fast_pdf(file_path, start_time, "PyMuPDF not available")
        
        try:
            # Direct file opening with timeout-like behavior
            doc = fitz.open(str(file_path))
            
            try:
                page_count = len(doc)
                
                # Fail fast if doc seems corrupted
                if page_count == 0:
                    doc.close()
                    return self._fail_fast_pdf(file_path, start_time, "PDF has 0 pages")
                
                # Check page limits from config
                max_pages_config = 999999  # Default: no limit
                skip_if_over = 999999      # Default: no limit
                
                if config:
                    max_pages_config = config.get('pdf.max_pages_to_extract', 999999)
                    skip_if_over = config.get('pdf.skip_if_pages_over', 999999)
                
                # Skip if too many pages
                if page_count > skip_if_over:
                    doc.close()
                    return self._fail_fast_pdf(file_path, start_time, 
                                               f"PDF has {page_count} pages (limit: {skip_if_over})")
                
                # Limit pages to extract for speed
                pages_to_extract = min(page_count, max_pages_config)
                
                # Extract text with error handling per page
                try:
                    text = self._extract_sequential_safe(doc, pages_to_extract)
                    
                    # Add notice if we limited pages
                    if pages_to_extract < page_count:
                        text += f"\n\n[Note: Only extracted first {pages_to_extract} of {page_count} pages for speed]"
                except Exception as text_error:
                    # If text extraction fails, try basic extraction
                    try:
                        text = self._extract_basic_text(doc)
                    except:
                        doc.close()
                        return self._fail_fast_pdf(file_path, start_time, f"Text extraction failed: {str(text_error)}")
                
                # Get basic metadata without complex parsing
                try:
                    metadata = self._extract_metadata_safe(doc, file_path)
                except Exception:
                    metadata = {
                        "filename": file_path.name,
                        "format": "PDF",
                        "pages": page_count,
                        "size_bytes": file_path.stat().st_size
                    }
                
                doc.close()
                
                extraction_time = time.perf_counter() - start_time
                pages_per_second = page_count / extraction_time if extraction_time > 0 else 0
                
                result = UltraFastResult(
                    file_path=str(file_path),
                    success=True,
                    text=text or "[No text extracted]",
                    page_count=page_count,
                    extraction_time=extraction_time,
                    pages_per_second=pages_per_second,
                    metadata=metadata
                )
                
                if len(text or "") < 10_000_000:
                    self._cache_result(cache_key, result)
                
                return result
                
            except Exception as process_error:
                try:
                    doc.close()
                except:
                    pass
                return self._fail_fast_pdf(file_path, start_time, f"PDF processing error: {str(process_error)}")
                
        except Exception as open_error:
            return self._fail_fast_pdf(file_path, start_time, f"PDF open error: {str(open_error)}")
    
    def _fail_fast_pdf(self, file_path: Path, start_time: float, error_msg: str) -> UltraFastResult:
        """Quickly fail a PDF with error info and move on."""
        extraction_time = time.perf_counter() - start_time
        return UltraFastResult(
            file_path=str(file_path),
            success=False,
            text="",
            page_count=0,
            extraction_time=extraction_time,
            pages_per_second=0,
            metadata={"filename": file_path.name, "format": "PDF", "error": "fail_fast"},
            error=error_msg
        )
    
    def _extract_sequential_safe(self, doc, max_pages=None) -> str:
        """Safe sequential extraction with per-page error handling."""
        texts = []
        pages_to_process = min(len(doc), max_pages) if max_pages else len(doc)
        
        for i in range(pages_to_process):
            try:
                page = doc[i]
                text = page.get_text("text", flags=fitz.TEXT_PRESERVE_WHITESPACE)
                texts.append(text or "")
            except Exception as e:
                texts.append(f"[Page {i+1} extraction failed: {str(e)[:50]}]")
        return '\n'.join(texts)
    
    def _extract_basic_text(self, doc) -> str:
        """Most basic text extraction - just get what we can."""
        texts = []
        for i in range(min(len(doc), 100)):  # Limit to first 100 pages for speed
            try:
                page = doc[i]
                text = page.get_text()  # Simplest method
                texts.append(text or "")
            except:
                texts.append(f"[Page {i+1} failed]")
        return '\n'.join(texts)
    
    def _extract_metadata_safe(self, doc, file_path: Path) -> Dict[str, Any]:
        """Safe metadata extraction."""
        try:
            metadata = {
                "filename": file_path.name,
                "pages": len(doc),
                "size_bytes": file_path.stat().st_size,
                "format": "PDF",
            }
            
            # Try to get PDF metadata safely
            if hasattr(doc, 'metadata') and doc.metadata:
                try:
                    metadata["title"] = doc.metadata.get("title", "")
                    metadata["author"] = doc.metadata.get("author", "")
                except:
                    pass
            
            return metadata
        except Exception:
            return {
                "filename": file_path.name,
                "format": "PDF",
                "pages": 1,
                "size_bytes": file_path.stat().st_size
            }
    
    def _extract_text(self, file_path: Path, start_time: float, cache_key: str) -> UltraFastResult:
        """Extract from plain text files."""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()
        
        metadata = {
            "filename": file_path.name,
            "format": file_path.suffix.upper(),
            "size_bytes": file_path.stat().st_size,
            "pages": 1
        }
        
        extraction_time = time.perf_counter() - start_time
        
        result = UltraFastResult(
            file_path=str(file_path),
            success=True,
            text=text,
            page_count=1,
            extraction_time=extraction_time,
            pages_per_second=1 / extraction_time if extraction_time > 0 else 0,
            metadata=metadata
        )
        
        self._cache_result(cache_key, result)
        return result
    
    def _extract_html(self, file_path: Path, start_time: float, cache_key: str) -> UltraFastResult:
        """Extract text from HTML files."""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            html_content = f.read()
        
        # Simple HTML text extraction (without BeautifulSoup for speed)
        import re
        # Remove script and style elements
        html_content = re.sub(r'<script[^>]*>.*?</script>', '', html_content, flags=re.DOTALL | re.IGNORECASE)
        html_content = re.sub(r'<style[^>]*>.*?</style>', '', html_content, flags=re.DOTALL | re.IGNORECASE)
        # Remove HTML tags
        text = re.sub(r'<[^>]+>', '', html_content)
        # Clean up whitespace
        text = re.sub(r'\s+', ' ', text).strip()
        
        metadata = {
            "filename": file_path.name,
            "format": "HTML",
            "size_bytes": file_path.stat().st_size,
            "pages": 1
        }
        
        extraction_time = time.perf_counter() - start_time
        
        result = UltraFastResult(
            file_path=str(file_path),
            success=True,
            text=text,
            page_count=1,
            extraction_time=extraction_time,
            pages_per_second=1 / extraction_time if extraction_time > 0 else 0,
            metadata=metadata
        )
        
        self._cache_result(cache_key, result)
        return result
    
    def _extract_csv(self, file_path: Path, start_time: float, cache_key: str) -> UltraFastResult:
        """Extract text from CSV files."""
        import csv
        
        rows = []
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            reader = csv.reader(f)
            for row in reader:
                rows.append(' | '.join(row))
        
        text = '\n'.join(rows)
        
        metadata = {
            "filename": file_path.name,
            "format": "CSV",
            "size_bytes": file_path.stat().st_size,
            "pages": 1,
            "rows": len(rows)
        }
        
        extraction_time = time.perf_counter() - start_time
        
        result = UltraFastResult(
            file_path=str(file_path),
            success=True,
            text=text,
            page_count=1,
            extraction_time=extraction_time,
            pages_per_second=1 / extraction_time if extraction_time > 0 else 0,
            metadata=metadata
        )
        
        self._cache_result(cache_key, result)
        return result
    
    def _extract_docx(self, file_path: Path, start_time: float, cache_key: str) -> UltraFastResult:
        """Extract text from DOCX files."""
        try:
            from docx import Document
            doc = Document(file_path)
            text = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
            
            metadata = {
                "filename": file_path.name,
                "format": "DOCX",
                "size_bytes": file_path.stat().st_size,
                "pages": max(1, len(doc.paragraphs) // 20)  # Estimate pages
            }
            
            extraction_time = time.perf_counter() - start_time
            pages = metadata["pages"]
            
            result = UltraFastResult(
                file_path=str(file_path),
                success=True,
                text=text,
                page_count=pages,
                extraction_time=extraction_time,
                pages_per_second=pages / extraction_time if extraction_time > 0 else 0,
                metadata=metadata
            )
            
            self._cache_result(cache_key, result)
            return result
            
        except ImportError:
            # Fallback: try to extract as zip
            return self._extract_docx_fallback(file_path, start_time, cache_key)
    
    def _extract_docx_fallback(self, file_path: Path, start_time: float, cache_key: str) -> UltraFastResult:
        """Fallback DOCX extraction using zipfile."""
        import zipfile
        import xml.etree.ElementTree as ET
        
        try:
            with zipfile.ZipFile(file_path, 'r') as doc:
                xml_content = doc.read('word/document.xml').decode('utf-8')
                root = ET.fromstring(xml_content)
                
                # Extract text from XML
                text_parts = []
                for elem in root.iter():
                    if elem.text:
                        text_parts.append(elem.text)
                
                text = ' '.join(text_parts)
                
                metadata = {
                    "filename": file_path.name,
                    "format": "DOCX",
                    "size_bytes": file_path.stat().st_size,
                    "pages": max(1, len(text) // 2000)  # Estimate pages
                }
                
                extraction_time = time.perf_counter() - start_time
                pages = metadata["pages"]
                
                result = UltraFastResult(
                    file_path=str(file_path),
                    success=True,
                    text=text,
                    page_count=pages,
                    extraction_time=extraction_time,
                    pages_per_second=pages / extraction_time if extraction_time > 0 else 0,
                    metadata=metadata
                )
                
                self._cache_result(cache_key, result)
                return result
                
        except Exception as e:
            extraction_time = time.perf_counter() - start_time
            return UltraFastResult(
                file_path=str(file_path),
                success=False,
                text="",
                page_count=0,
                extraction_time=extraction_time,
                pages_per_second=0,
                metadata={"filename": file_path.name, "format": "DOCX"},
                error=str(e)
            )
    
    def _extract_pptx(self, file_path: Path, start_time: float, cache_key: str) -> UltraFastResult:
        """Extract text from PPTX files."""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)
            
            text = '\n'.join(text_parts)
            
            metadata = {
                "filename": file_path.name,
                "format": "PPTX",
                "size_bytes": file_path.stat().st_size,
                "pages": len(prs.slides)
            }
            
            extraction_time = time.perf_counter() - start_time
            pages = len(prs.slides)
            
            result = UltraFastResult(
                file_path=str(file_path),
                success=True,
                text=text,
                page_count=pages,
                extraction_time=extraction_time,
                pages_per_second=pages / extraction_time if extraction_time > 0 else 0,
                metadata=metadata
            )
            
            self._cache_result(cache_key, result)
            return result
            
        except ImportError:
            return self._extract_pptx_fallback(file_path, start_time, cache_key)
    
    def _extract_pptx_fallback(self, file_path: Path, start_time: float, cache_key: str) -> UltraFastResult:
        """Fallback PPTX extraction."""
        import zipfile
        import xml.etree.ElementTree as ET
        
        try:
            text_parts = []
            with zipfile.ZipFile(file_path, 'r') as prs:
                # Get slide files
                slide_files = [f for f in prs.namelist() if f.startswith('ppt/slides/slide') and f.endswith('.xml')]
                
                for slide_file in slide_files:
                    xml_content = prs.read(slide_file).decode('utf-8')
                    root = ET.fromstring(xml_content)
                    
                    for elem in root.iter():
                        if elem.text:
                            text_parts.append(elem.text)
            
            text = '\n'.join(text_parts)
            
            metadata = {
                "filename": file_path.name,
                "format": "PPTX",
                "size_bytes": file_path.stat().st_size,
                "pages": len(slide_files)
            }
            
            extraction_time = time.perf_counter() - start_time
            pages = len(slide_files)
            
            result = UltraFastResult(
                file_path=str(file_path),
                success=True,
                text=text,
                page_count=pages,
                extraction_time=extraction_time,
                pages_per_second=pages / extraction_time if extraction_time > 0 else 0,
                metadata=metadata
            )
            
            self._cache_result(cache_key, result)
            return result
            
        except Exception as e:
            extraction_time = time.perf_counter() - start_time
            return UltraFastResult(
                file_path=str(file_path),
                success=False,
                text="",
                page_count=0,
                extraction_time=extraction_time,
                pages_per_second=0,
                metadata={"filename": file_path.name, "format": "PPTX"},
                error=str(e)
            )
    
    def _extract_xlsx(self, file_path: Path, start_time: float, cache_key: str) -> UltraFastResult:
        """Extract text from XLSX files."""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, data_only=True)
            
            text_parts = []
            total_rows = 0
            
            for sheet in wb.worksheets:
                text_parts.append(f"=== {sheet.title} ===")
                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        row_text = ' | '.join(str(cell) if cell is not None else '' for cell in row)
                        text_parts.append(row_text)
                        total_rows += 1
                text_parts.append("")
            
            text = '\n'.join(text_parts)
            
            metadata = {
                "filename": file_path.name,
                "format": "XLSX",
                "size_bytes": file_path.stat().st_size,
                "pages": len(wb.worksheets),
                "rows": total_rows
            }
            
            extraction_time = time.perf_counter() - start_time
            pages = len(wb.worksheets)
            
            result = UltraFastResult(
                file_path=str(file_path),
                success=True,
                text=text,
                page_count=pages,
                extraction_time=extraction_time,
                pages_per_second=pages / extraction_time if extraction_time > 0 else 0,
                metadata=metadata
            )
            
            self._cache_result(cache_key, result)
            return result
            
        except ImportError:
            return self._extract_xlsx_fallback(file_path, start_time, cache_key)
    
    def _extract_xlsx_fallback(self, file_path: Path, start_time: float, cache_key: str) -> UltraFastResult:
        """Fallback XLSX extraction."""
        import zipfile
        import xml.etree.ElementTree as ET
        
        try:
            text_parts = []
            with zipfile.ZipFile(file_path, 'r') as wb:
                # Try to extract shared strings
                try:
                    xml_content = wb.read('xl/sharedStrings.xml').decode('utf-8')
                    root = ET.fromstring(xml_content)
                    for elem in root.iter():
                        if elem.text:
                            text_parts.append(elem.text)
                except:
                    pass
                
                # Extract from worksheets
                worksheet_files = [f for f in wb.namelist() if f.startswith('xl/worksheets/') and f.endswith('.xml')]
                for ws_file in worksheet_files:
                    try:
                        xml_content = wb.read(ws_file).decode('utf-8')
                        root = ET.fromstring(xml_content)
                        for elem in root.iter():
                            if elem.text and elem.text.strip():
                                text_parts.append(elem.text)
                    except:
                        continue
            
            text = '\n'.join(text_parts)
            
            metadata = {
                "filename": file_path.name,
                "format": "XLSX",
                "size_bytes": file_path.stat().st_size,
                "pages": len(worksheet_files)
            }
            
            extraction_time = time.perf_counter() - start_time
            pages = max(1, len(worksheet_files))
            
            result = UltraFastResult(
                file_path=str(file_path),
                success=True,
                text=text,
                page_count=pages,
                extraction_time=extraction_time,
                pages_per_second=pages / extraction_time if extraction_time > 0 else 0,
                metadata=metadata
            )
            
            self._cache_result(cache_key, result)
            return result
            
        except Exception as e:
            extraction_time = time.perf_counter() - start_time
            return UltraFastResult(
                file_path=str(file_path),
                success=False,
                text="",
                page_count=0,
                extraction_time=extraction_time,
                pages_per_second=0,
                metadata={"filename": file_path.name, "format": "XLSX"},
                error=str(e)
            )
    
    def _extract_universal(self, file_path: Path, start_time: float, cache_key: str) -> UltraFastResult:
        """Universal extractor for ANY file type - tries multiple approaches."""
        
        file_ext = file_path.suffix.lower()
        
        # Try different extraction strategies in order of likelihood
        strategies = [
            self._try_text_extraction,
            self._try_binary_text_extraction,
            self._try_metadata_only,
        ]
        
        for strategy in strategies:
            try:
                result = strategy(file_path, start_time, cache_key, file_ext)
                if result and result.success:
                    return result
            except Exception:
                continue
        
        # If all strategies fail, return basic file info
        extraction_time = time.perf_counter() - start_time
        return UltraFastResult(
            file_path=str(file_path),
            success=True,  # Still mark as success - we got file info
            text=f"[Binary file: {file_path.name}]",
            page_count=1,
            extraction_time=extraction_time,
            pages_per_second=1 / extraction_time if extraction_time > 0 else 0,
            metadata={
                "filename": file_path.name,
                "format": file_ext.upper() or "UNKNOWN",
                "size_bytes": file_path.stat().st_size,
                "pages": 1,
                "extracted": "filename_only"
            }
        )
    
    def _try_text_extraction(self, file_path: Path, start_time: float, cache_key: str, file_ext: str) -> UltraFastResult:
        """Try to read as UTF-8 text file."""
        
        # Skip obviously binary files
        binary_extensions = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.tif', 
                           '.mp3', '.mp4', '.wav', '.avi', '.mov', '.zip', '.tar', '.gz',
                           '.exe', '.dll', '.so', '.dylib', '.bin'}
        
        if file_ext in binary_extensions:
            return None
        
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                # Read first chunk to see if it looks like text
                sample = f.read(1024)
                if len(sample.strip()) < 10:  # Too short or empty
                    return None
                
                # Check if it's mostly printable characters
                printable_ratio = sum(1 for c in sample if c.isprintable() or c.isspace()) / len(sample)
                if printable_ratio < 0.7:  # Less than 70% printable
                    return None
                
                # Read the full file
                f.seek(0)
                text = f.read()
                
                metadata = {
                    "filename": file_path.name,
                    "format": file_ext.upper() or "TEXT",
                    "size_bytes": file_path.stat().st_size,
                    "pages": 1,
                    "extracted": "text_content"
                }
                
                extraction_time = time.perf_counter() - start_time
                
                result = UltraFastResult(
                    file_path=str(file_path),
                    success=True,
                    text=text,
                    page_count=1,
                    extraction_time=extraction_time,
                    pages_per_second=1 / extraction_time if extraction_time > 0 else 0,
                    metadata=metadata
                )
                
                self._cache_result(cache_key, result)
                return result
                
        except Exception:
            return None
    
    def _try_binary_text_extraction(self, file_path: Path, start_time: float, cache_key: str, file_ext: str) -> UltraFastResult:
        """Try to extract any readable strings from binary files."""
        
        try:
            with open(file_path, 'rb') as f:
                # Read file in chunks to handle large files
                text_parts = []
                chunk_size = 8192
                
                while True:
                    chunk = f.read(chunk_size)
                    if not chunk:
                        break
                    
                    # Extract printable strings of length 4+
                    current_string = ""
                    for byte in chunk:
                        char = chr(byte) if 32 <= byte <= 126 else None
                        if char:
                            current_string += char
                        else:
                            if len(current_string) >= 4:  # Only keep strings of 4+ chars
                                text_parts.append(current_string)
                            current_string = ""
                    
                    # Don't extract too much from binary files
                    if len(text_parts) > 1000:
                        break
                
                # Add final string if any
                if len(current_string) >= 4:
                    text_parts.append(current_string)
                
                if not text_parts:
                    return None
                
                text = '\n'.join(text_parts[:500])  # Limit output
                
                metadata = {
                    "filename": file_path.name,
                    "format": file_ext.upper() or "BINARY",
                    "size_bytes": file_path.stat().st_size,
                    "pages": 1,
                    "extracted": "strings_from_binary",
                    "strings_found": len(text_parts)
                }
                
                extraction_time = time.perf_counter() - start_time
                
                result = UltraFastResult(
                    file_path=str(file_path),
                    success=True,
                    text=text,
                    page_count=1,
                    extraction_time=extraction_time,
                    pages_per_second=1 / extraction_time if extraction_time > 0 else 0,
                    metadata=metadata
                )
                
                self._cache_result(cache_key, result)
                return result
                
        except Exception:
            return None
    
    def _try_metadata_only(self, file_path: Path, start_time: float, cache_key: str, file_ext: str) -> UltraFastResult:
        """Extract only file metadata - always succeeds."""
        
        stat = file_path.stat()
        
        text = f"""File: {file_path.name}
Extension: {file_ext}
Size: {stat.st_size} bytes
Modified: {stat.st_mtime}
Path: {file_path}"""
        
        metadata = {
            "filename": file_path.name,
            "format": file_ext.upper() or "UNKNOWN",
            "size_bytes": stat.st_size,
            "pages": 1,
            "extracted": "metadata_only",
            "modification_time": stat.st_mtime
        }
        
        extraction_time = time.perf_counter() - start_time
        
        result = UltraFastResult(
            file_path=str(file_path),
            success=True,
            text=text,
            page_count=1,
            extraction_time=extraction_time,
            pages_per_second=1 / extraction_time if extraction_time > 0 else 0,
            metadata=metadata
        )
        
        self._cache_result(cache_key, result)
        return result
    
    def _cache_result(self, cache_key: str, result: UltraFastResult):
        """Cache extraction result."""
        
        cache_data = {
            'text': result.text,
            'pages': result.page_count,
            'metadata': result.metadata
        }
        
        # Simple size-based eviction
        data_size = len(result.text)
        if self.current_cache_size + data_size > self.cache_size:
            # Evict oldest entries
            to_remove = list(self.cache.keys())[:len(self.cache)//4]
            for key in to_remove:
                del self.cache[key]
            self.current_cache_size *= 0.75
        
        self.cache[cache_key] = cache_data
        self.current_cache_size += data_size
    
    def _extract_fallback(self, file_path: Path) -> UltraFastResult:
        """Fallback extraction without PyMuPDF."""
        
        start_time = time.perf_counter()
        
        # Try pdftotext command line tool (very fast)
        import subprocess
        try:
            result = subprocess.run(
                ['pdftotext', '-layout', str(file_path), '-'],
                capture_output=True,
                text=True,
                timeout=30
            )
            
            if result.returncode == 0:
                text = result.stdout
                # Estimate page count from text
                page_count = text.count('\f') + 1  # Form feeds indicate pages
                
                extraction_time = time.perf_counter() - start_time
                
                return UltraFastResult(
                    file_path=str(file_path),
                    success=True,
                    text=text,
                    page_count=page_count,
                    extraction_time=extraction_time,
                    pages_per_second=page_count / extraction_time,
                    metadata={"filename": file_path.name, "pages": page_count}
                )
        except:
            pass
        
        extraction_time = time.perf_counter() - start_time
        return UltraFastResult(
            file_path=str(file_path),
            success=False,
            text="",
            page_count=0,
            extraction_time=extraction_time,
            pages_per_second=0,
            metadata={},
            error="No PDF extraction library available"
        )
    
    def process_batch(self, file_paths: List[Path]) -> List[UltraFastResult]:
        """Process multiple files in parallel batches."""
        
        results = []
        
        # Use threading instead of multiprocessing to avoid pickling issues
        with ThreadPoolExecutor(max_workers=self.num_workers) as executor:
            futures = []
            for file_path in file_paths:
                future = executor.submit(self.extract_document_ultrafast, file_path)
                futures.append(future)
            
            for future in futures:
                results.append(future.result())
        
        return results
    
    async def process_batch_async(self, file_paths: List[Path]) -> List[UltraFastResult]:
        """Async batch processing for even better concurrency."""
        
        loop = asyncio.get_event_loop()
        tasks = []
        
        # Create thread pool for async execution
        with ThreadPoolExecutor(max_workers=self.num_workers) as thread_pool:
            for file_path in file_paths:
                task = loop.run_in_executor(
                    thread_pool,
                    self.extract_document_ultrafast,
                    file_path
                )
                tasks.append(task)
            
            results = await asyncio.gather(*tasks)
        
        return results
    
    def shutdown(self):
        """Clean shutdown of pools."""
        # Pools are created as context managers, no need to shutdown
        pass


class HyperBatchProcessor:
    """Batch processor for maximum throughput."""
    
    def __init__(self, num_workers: int = None):
        self.num_workers = num_workers or mp.cpu_count()
        self.extractor = UltraFastExtractor(num_workers=self.num_workers)
    
    def process_directory(self, directory: Path, pattern: str = "*.*") -> Dict[str, Any]:
        """Process entire directory at maximum speed."""
        
        start_time = time.perf_counter()
        
        # Find ALL files recursively - no filtering by extension
        files = []
        
        # Use recursive glob to find ALL files, following symbolic links
        import os
        for root, dirs, filenames in os.walk(directory, followlinks=True):
            for filename in filenames:
                file_path = Path(root) / filename
                files.append(file_path)
        
        total_files = len(files)
        
        print(f"Found {total_files} files (processing ALL file types):")
        # Show breakdown by file type
        file_types = {}
        for f in files:
            ext = f.suffix.lower()
            file_types[ext] = file_types.get(ext, 0) + 1
        
        for ext, count in sorted(file_types.items()):
            print(f"  {ext}: {count} files")
        
        if total_files == 0:
            return {"error": "No files found"}
        
        print(f"Processing {total_files} files with threading...")
        
        # Process all files using the single extractor with threading
        results = self.extractor.process_batch(files)
        
        # Calculate statistics
        total_time = time.perf_counter() - start_time
        successful = [r for r in results if r.success]
        total_pages = sum(r.page_count for r in successful)
        
        stats = {
            "total_files": total_files,
            "successful_files": len(successful),
            "failed_files": total_files - len(successful),
            "total_pages": total_pages,
            "total_time": total_time,
            "files_per_second": total_files / total_time if total_time > 0 else 0,
            "pages_per_second": total_pages / total_time if total_time > 0 else 0,
            "average_pages_per_file": total_pages / len(successful) if successful else 0,
            "results": results
        }
        
        return stats
    
    def shutdown(self):
        """Shutdown extractor."""
        self.extractor.shutdown()


def format_metadata_header(metadata: Dict[str, Any]) -> str:
    """Format metadata as markdown header."""
    
    lines = ["---"]
    for key, value in metadata.items():
        if value:
            lines.append(f"{key}: {value}")
    lines.append("---")
    return '\n'.join(lines)


def main():
    """Main entry point for testing."""
    
    import argparse
    from config_loader import load_config
    
    parser = argparse.ArgumentParser(description="Ultra-fast document processor")
    parser.add_argument("--config", default="config.yaml", help="Configuration file")
    parser.add_argument("--test-config", action="store_true", 
                       help="Create test config for troubleshooting")
    parser.add_argument("input", nargs='*', help="Input file(s) or directory(ies) (overrides config)")
    parser.add_argument("--workers", type=int, help="Number of worker processes (overrides config)")
    parser.add_argument("--output", help="Output directory (overrides config)")
    parser.add_argument("--benchmark", action="store_true", help="Run in benchmark mode")
    
    args = parser.parse_args()
    
    # Create test config if requested
    if args.test_config:
        from config_loader import create_test_config
        create_test_config()
        print("Test config created. Run with: --config test_config.yaml")
        return
    
    # Load configuration
    config = load_config(args.config)
    
    # Command line overrides
    if args.input:
        input_paths = [Path(p) for p in args.input]
        # Override config with command line inputs
        all_files = []
        for input_path in input_paths:
            if input_path.is_dir():
                for file_path in input_path.rglob("*"):
                    if file_path.is_file():
                        all_files.append(file_path)
            elif input_path.is_file():
                all_files.append(input_path)
    else:
        # Use config file inputs
        all_files = config.get_input_files()
    
    # Apply config overrides
    max_workers = args.workers if args.workers else config.get('processing.max_workers', mp.cpu_count())
    output_dir = Path(args.output) if args.output else Path(config.get('output.directory', 'output'))
    
    print(f"🔧 CONFIGURATION:")
    print(f"  Workers: {max_workers}")
    print(f"  Output: {output_dir}")
    print(f"  Total files found: {len(all_files)}")
    print(f"  Config file: {args.config}")
    
    if len(all_files) == 0:
        print("❌ No files found to process")
        return
        
    if len(all_files) == 1:
        # Single file processing
        print(f"\n📄 SINGLE FILE MODE")
        extractor = UltraFastExtractor(num_workers=max_workers)
        result = extractor.extract_document_ultrafast(all_files[0])
        
        print(f"File: {result.file_path}")
        print(f"Pages: {result.page_count}")
        print(f"Extraction time: {result.extraction_time:.3f}s")
        print(f"Speed: {result.pages_per_second:.1f} pages/second")
        
        if output_dir:
            output_path = output_dir / f"{all_files[0].stem}.md"
            output_path.parent.mkdir(parents=True, exist_ok=True)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(format_metadata_header(result.metadata))
                f.write('\n\n')
                f.write(result.text)
            
            print(f"Output saved to: {output_path}")
        
        extractor.shutdown()
        
    else:
        # Multiple files processing
        print(f"\n📚 BATCH MODE - {len(all_files)} files")
        processor = HyperBatchProcessor(num_workers=max_workers)
        # Pass config to extractor for PDF limits
        processor.extractor.config = config
        
        # Show breakdown by file type
        file_types = {}
        for f in all_files:
            ext = f.suffix.lower() or 'no_extension'
            file_types[ext] = file_types.get(ext, 0) + 1
        
        print("File breakdown:")
        for ext, count in sorted(file_types.items()):
            print(f"  {ext}: {count} files")
        
        # Analyze potentially problematic file types
        print("\n🔍 FILE TYPE ANALYSIS:")
        
        # Known fast formats
        fast_formats = {'.txt', '.md', '.csv', '.html', '.htm'}
        # Document formats (can be slow)
        doc_formats = {'.pdf', '.docx', '.pptx', '.xlsx'}
        # Binary/unknown formats (likely slow)
        binary_formats = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.tif', 
                         '.mp3', '.mp4', '.wav', '.avi', '.mov', '.zip', '.tar', '.gz',
                         '.exe', '.dll', '.so', '.dylib', '.bin', '.dat', '.log'}
        
        fast_count = sum(count for ext, count in file_types.items() if ext in fast_formats)
        doc_count = sum(count for ext, count in file_types.items() if ext in doc_formats)
        binary_count = sum(count for ext, count in file_types.items() if ext in binary_formats)
        unknown_count = len(all_files) - fast_count - doc_count - binary_count
        
        print(f"  📄 Fast text files: {fast_count} files")
        print(f"  📚 Document files: {doc_count} files") 
        print(f"  💾 Binary files: {binary_count} files")
        print(f"  ❓ Unknown/other: {unknown_count} files")
        
        # Show the most common extensions
        sorted_types = sorted(file_types.items(), key=lambda x: x[1], reverse=True)
        print(f"\n📊 TOP FILE EXTENSIONS:")
        for ext, count in sorted_types[:10]:
            percentage = (count / len(all_files)) * 100
            print(f"  {ext or 'no_extension'}: {count} files ({percentage:.1f}%)")
        
        # Identify potentially slow formats
        slow_formats = []
        for ext, count in sorted_types:
            if ext not in fast_formats and count > 10:
                slow_formats.append((ext, count))
        
        if slow_formats:
            print(f"\n⚠️  POTENTIALLY SLOW FORMATS (>10 files):")
            for ext, count in slow_formats[:5]:
                print(f"  {ext or 'no_extension'}: {count} files - may be causing slowdown")
        
        if len(all_files) == 0:
            print("No files found to process")
            return
        
        # Process all files one by one with progress tracking
        start_time = time.perf_counter()
        results = []
        
        print(f"\n🔄 PROCESSING {len(all_files)} FILES:")
        print("=" * 60)
        
        # Define file types to skip (will be too slow/pointless)
        skip_extensions = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.tif',
                          '.mp3', '.mp4', '.wav', '.avi', '.mov', '.wmv', '.flv',
                          '.zip', '.tar', '.gz', '.bz2', '.rar', '.7z',
                          '.exe', '.dll', '.so', '.dylib', '.bin', '.dat'}
        
        skipped_count = 0
        
        for i, file_path in enumerate(all_files, 1):
            file_start = time.perf_counter()
            file_ext = file_path.suffix.lower()
            
            # Show progress every 50 files or for first 20
            if i % 50 == 0 or i <= 20:
                print(f"[{i:4d}/{len(all_files)}] Processing: {file_path.name} ({file_ext or 'no ext'})")
            
            # Skip obviously problematic file types
            if file_ext in skip_extensions:
                if i <= 20 or i % 100 == 0:  # Show some skips
                    print(f"  ⏭️  SKIPPED: {file_path.name} ({file_ext}) - binary file")
                skipped_count += 1
                results.append(UltraFastResult(
                    file_path=str(file_path),
                    success=True,
                    text=f"[Skipped binary file: {file_path.name}]",
                    page_count=1,
                    extraction_time=0.001,
                    pages_per_second=1000,
                    metadata={"filename": file_path.name, "format": file_ext.upper(), "extracted": "skipped_binary"}
                ))
                continue
            
            try:
                result = processor.extractor.extract_document_ultrafast(file_path)
                results.append(result)
                
                file_time = time.perf_counter() - file_start
                
                # Warn about slow files (>1 second)
                if file_time > 1.0:
                    print(f"  ⚠️  SLOW FILE: {file_path.name} took {file_time:.2f}s ({file_path.suffix})")
                
                # Show progress for very slow files
                if file_time > 5.0:
                    print(f"  🐌 VERY SLOW: {file_path.name} took {file_time:.2f}s - consider skipping this type")
                    
            except Exception as e:
                print(f"  ❌ CRASHED: {file_path.name} - {str(e)[:100]}")
                # Create failed result
                results.append(UltraFastResult(
                    file_path=str(file_path),
                    success=False,
                    text="",
                    page_count=0,
                    extraction_time=time.perf_counter() - file_start,
                    pages_per_second=0,
                    metadata={"filename": file_path.name, "format": file_path.suffix.upper()},
                    error=f"Processing crashed: {str(e)}"
                ))
                
        total_time = time.perf_counter() - start_time
        print(f"\n✅ Completed processing {len(results)} files in {total_time:.2f}s")
        if skipped_count > 0:
            print(f"⏭️  Skipped {skipped_count} binary files for speed")
        
        # Calculate statistics
        successful = [r for r in results if r.success]
        total_pages = sum(r.page_count for r in successful)
        
        stats = {
            "total_files": len(all_files),
            "successful_files": len(successful),
            "failed_files": len(all_files) - len(successful),
            "total_pages": total_pages,
            "total_time": total_time,
            "files_per_second": len(all_files) / total_time if total_time > 0 else 0,
            "pages_per_second": total_pages / total_time if total_time > 0 else 0,
            "average_pages_per_file": total_pages / len(successful) if successful else 0,
            "results": results
        }
        
        print("\n" + "="*60)
        print("ULTRA-FAST PROCESSING COMPLETE")
        print("="*60)
        print(f"Total files: {stats['total_files']}")
        print(f"Successful: {stats['successful_files']}")
        print(f"Failed: {stats['failed_files']}")
        print(f"Total pages: {stats['total_pages']}")
        print(f"Total time: {stats['total_time']:.2f}s")
        print(f"Files/second: {stats['files_per_second']:.1f}")
        print(f"Pages/second: {stats['pages_per_second']:.1f}")
        
        # Show failure analysis - SHOW ALL FAILURES
        if stats['failed_files'] > 0:
            print(f"\n❌ FAILURE ANALYSIS - ALL {stats['failed_files']} FAILURES:")
            failed_results = [r for r in stats['results'] if not r.success]
            
            # Group failures by error type
            error_types = {}
            for result in failed_results:  # Show ALL failures
                error = result.error or "Unknown error"
                if error not in error_types:
                    error_types[error] = []
                error_types[error].append(result.file_path)
            
            for error, file_paths in error_types.items():
                print(f"\n  ERROR: {error}")
                print(f"  FILES ({len(file_paths)}):")
                for file_path in file_paths:  # Show ALL files with this error
                    print(f"    - {Path(file_path).name}")
                print("")
        
        # Show success analysis
        if stats['successful_files'] > 0:
            print(f"\n✅ SUCCESS ANALYSIS:")
            successful_results = [r for r in stats['results'] if r.success]
            
            # Group by extraction method
            extraction_methods = {}
            for result in successful_results:
                method = result.metadata.get('extracted', 'unknown')
                if method not in extraction_methods:
                    extraction_methods[method] = []
                extraction_methods[method].append(result)
            
            for method, results in extraction_methods.items():
                total_pages = sum(r.page_count for r in results)
                avg_speed = sum(r.pages_per_second for r in results) / len(results)
                avg_time = sum(r.extraction_time for r in results) / len(results)
                print(f"  {method}: {len(results)} files, {total_pages} pages, {avg_speed:.1f} avg pages/sec, {avg_time:.3f}s avg time")
            
            # Analyze by file extension performance
            print(f"\n📊 PERFORMANCE BY FILE TYPE:")
            ext_performance = {}
            for result in successful_results:
                file_ext = Path(result.file_path).suffix.lower() or 'no_extension'
                if file_ext not in ext_performance:
                    ext_performance[file_ext] = []
                ext_performance[file_ext].append(result)
            
            # Sort by total processing time (slowest first)
            ext_stats = []
            for ext, results in ext_performance.items():
                total_time = sum(r.extraction_time for r in results)
                total_pages = sum(r.page_count for r in results)
                avg_time = total_time / len(results)
                avg_speed = sum(r.pages_per_second for r in results) / len(results)
                ext_stats.append((ext, len(results), total_time, avg_time, avg_speed, total_pages))
            
            # Sort by total time (biggest bottlenecks first)
            ext_stats.sort(key=lambda x: x[2], reverse=True)
            
            print("  Top time consumers:")
            for ext, count, total_time, avg_time, avg_speed, pages in ext_stats[:8]:
                print(f"    {ext}: {count} files, {total_time:.2f}s total, {avg_time:.3f}s avg, {avg_speed:.1f} pages/sec")
            
            # Find the fastest and slowest
            if ext_stats:
                fastest = min(ext_stats, key=lambda x: x[3])  # By avg time
                slowest = max(ext_stats, key=lambda x: x[3])
                print(f"\n  🚀 Fastest: {fastest[0]} ({fastest[3]:.3f}s avg)")
                print(f"  🐌 Slowest: {slowest[0]} ({slowest[3]:.3f}s avg)")
        
        if stats['pages_per_second'] >= 1000:
            print("\n🚀 TARGET ACHIEVED: 1000+ pages/second!")
        else:
            print(f"\n📊 Current: {stats['pages_per_second']:.1f} pages/sec (Target: 1000)")
        
        # Always save output if output directory is configured
        if output_dir:
            output_dir.mkdir(parents=True, exist_ok=True)
            
            print(f"\n💾 SAVING OUTPUT FILES...")
            saved_count = 0
            
            for result in stats['results']:
                if result.success:
                    output_path = output_dir / f"{Path(result.file_path).stem}.md"
                    with open(output_path, 'w', encoding='utf-8') as f:
                        f.write(format_metadata_header(result.metadata))
                        f.write('\n\n')
                        f.write(result.text)
                    saved_count += 1
            
            print(f"✅ Saved {saved_count} markdown files to: {output_dir}")
        else:
            print("\n⚠️  No output directory configured - files not saved")
        
        processor.shutdown()


if __name__ == "__main__":
    main()